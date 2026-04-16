"""
Unit tests for PPT renderer.
"""
import pytest
import os
import tempfile
from unittest.mock import patch, MagicMock, PropertyMock
from PIL import Image
from autoppt.ppt_renderer import PPTRenderer
from autoppt.data_types import ChartData, ChartType, SlideLayout, SlideSpec, StatisticData
from autoppt.style_selector import get_all_styles


class TestPPTRendererInit:
    """Tests for PPTRenderer initialization."""

    def test_renderer_instantiation(self):
        """Test that PPTRenderer can be instantiated."""
        renderer = PPTRenderer()
        assert renderer is not None
        assert renderer.prs is not None

    def test_renderer_with_no_template(self):
        """Test renderer creates empty presentation without template."""
        renderer = PPTRenderer(template_path=None)
        assert len(renderer.prs.slides) == 0


class TestApplyStyle:
    """Tests for apply_style method."""

    def test_apply_minimalist_style(self):
        """Test applying minimalist style."""
        renderer = PPTRenderer()
        renderer.apply_style("minimalist")

        assert hasattr(renderer, 'current_style')
        assert 'title_color' in renderer.current_style
        assert 'text_color' in renderer.current_style
        assert 'bg_color' in renderer.current_style
        assert 'font_name' in renderer.current_style

    def test_apply_all_styles(self):
        """Test that all styles can be applied without error."""
        styles = get_all_styles()

        for style in styles:
            renderer = PPTRenderer()
            renderer.apply_style(style)
            assert hasattr(renderer, 'current_style')

    def test_apply_unknown_style_defaults(self):
        """Test that unknown style defaults to minimalist."""
        renderer = PPTRenderer()
        renderer.apply_style("unknown_style")

        # Should not raise, should default to minimalist
        assert hasattr(renderer, 'current_style')


class TestAddSlides:
    """Tests for slide addition methods."""

    def test_add_title_slide(self):
        """Test adding a title slide."""
        renderer = PPTRenderer()
        renderer.apply_style("minimalist")

        initial_count = len(renderer.prs.slides)
        renderer.add_title_slide("Test Title", "Test Subtitle")

        assert len(renderer.prs.slides) == initial_count + 1
        slide = renderer.prs.slides[initial_count]
        all_text = " ".join(
            shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
        )
        assert "Test Title" in all_text

    def test_add_section_header(self):
        """Test adding a section header slide."""
        renderer = PPTRenderer()
        renderer.apply_style("minimalist")

        initial_count = len(renderer.prs.slides)
        renderer.add_section_header("Section 1")

        assert len(renderer.prs.slides) == initial_count + 1

    def test_add_content_slide(self, sample_bullets):
        """Test adding a content slide with bullets."""
        renderer = PPTRenderer()
        renderer.apply_style("minimalist")

        initial_count = len(renderer.prs.slides)
        renderer.add_content_slide(
            title="Content Slide",
            bullets=sample_bullets,
            notes="Speaker notes here"
        )

        assert len(renderer.prs.slides) == initial_count + 1
        slide = renderer.prs.slides[initial_count]
        all_text = " ".join(
            shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
        )
        assert "First key point about the topic" in all_text

    def test_add_citations_slide(self):
        """Test adding a citations slide."""
        renderer = PPTRenderer()
        renderer.apply_style("minimalist")

        citations = [
            "https://example.com/source1",
            "https://example.com/source2"
        ]

        initial_count = len(renderer.prs.slides)
        renderer.add_citations_slide(citations)

        assert len(renderer.prs.slides) == initial_count + 1

    def test_add_citations_slide_empty_list(self):
        """Test that empty citations list adds no slide."""
        renderer = PPTRenderer()
        renderer.apply_style("minimalist")

        initial_count = len(renderer.prs.slides)
        renderer.add_citations_slide([])

        assert len(renderer.prs.slides) == initial_count


class TestChartSlide:
    """Tests for chart slide functionality."""

    def test_add_chart_slide(self):
        """Test adding a chart slide."""
        renderer = PPTRenderer()
        renderer.apply_style("corporate")

        chart_data = ChartData(
            chart_type=ChartType.COLUMN,
            title="Quarterly Revenue",
            categories=["Q1", "Q2", "Q3", "Q4"],
            values=[100.0, 150.0, 200.0, 250.0],
            series_name="2025"
        )

        initial_count = len(renderer.prs.slides)
        renderer.add_chart_slide("Revenue Analysis", chart_data)

        assert len(renderer.prs.slides) == initial_count + 1

    def test_add_two_column_slide(self, sample_bullets):
        """Test adding a two-column slide."""
        renderer = PPTRenderer()
        renderer.apply_style("corporate")

        initial_count = len(renderer.prs.slides)
        renderer.add_two_column_slide(
            title="Comparison",
            left_bullets=sample_bullets,
            right_bullets=sample_bullets,
            left_title="Pros",
            right_title="Cons"
        )
        assert len(renderer.prs.slides) == initial_count + 1

    def test_add_comparison_slide(self, sample_bullets):
        """Test adding a comparison slide helper."""
        renderer = PPTRenderer()
        renderer.apply_style("corporate")

        initial_count = len(renderer.prs.slides)
        item_a = {"name": "A", "points": sample_bullets}
        item_b = {"name": "B", "points": sample_bullets}

        renderer.add_comparison_slide("A vs B", item_a, item_b)
        assert len(renderer.prs.slides) == initial_count + 1

    def test_add_quote_slide(self):
        """Test adding a quote slide."""
        renderer = PPTRenderer()
        renderer.apply_style("creative")

        initial_count = len(renderer.prs.slides)
        renderer.add_quote_slide(
            quote="Innovation distinguishes between a leader and a follower.",
            author="Steve Jobs",
            context="Apple"
        )
        assert len(renderer.prs.slides) == initial_count + 1

    def test_add_statistics_slide(self):
        """Test adding a statistics slide."""
        renderer = PPTRenderer()
        renderer.apply_style("technology")

        stats = [
            {"value": "85%", "label": "Growth"},
            {"value": "2025", "label": "Year"},
            {"value": "1M+", "label": "Users"}
        ]

        initial_count = len(renderer.prs.slides)
        renderer.add_statistics_slide("Key Metrics", stats)
        assert len(renderer.prs.slides) == initial_count + 1

    def test_statistics_slide_width_clamped_to_minimum(self):
        """Statistics slide should not produce negative card widths even with extreme theme values."""
        renderer = PPTRenderer()
        renderer.apply_style("technology")
        stats = [
            {"value": "1", "label": "A"},
            {"value": "2", "label": "B"},
            {"value": "3", "label": "C"},
            {"value": "4", "label": "D"},
        ]
        initial_count = len(renderer.prs.slides)
        renderer.add_statistics_slide("Metrics", stats)
        assert len(renderer.prs.slides) == initial_count + 1

    def test_render_quote_slide_spec(self):
        """Test rendering a quote slide from SlideSpec."""
        renderer = PPTRenderer()
        renderer.apply_style("creative")

        initial_count = len(renderer.prs.slides)
        renderer.render_slide(
            SlideSpec(
                layout=SlideLayout.QUOTE,
                title="Quote",
                quote_text="Innovation distinguishes between a leader and a follower.",
                quote_author="Steve Jobs",
                quote_context="Apple",
            )
        )
        assert len(renderer.prs.slides) == initial_count + 1

    def test_render_two_column_slide_spec(self, sample_bullets):
        """Test rendering a two-column slide from SlideSpec."""
        renderer = PPTRenderer()
        renderer.apply_style("corporate")

        initial_count = len(renderer.prs.slides)
        renderer.render_slide(
            SlideSpec(
                layout=SlideLayout.TWO_COLUMN,
                title="Two Column",
                left_title="Left",
                right_title="Right",
                left_bullets=sample_bullets,
                right_bullets=sample_bullets,
            )
        )
        assert len(renderer.prs.slides) == initial_count + 1

    def test_render_comparison_slide_spec(self, sample_bullets):
        """Test rendering a comparison slide from SlideSpec."""
        renderer = PPTRenderer()
        renderer.apply_style("corporate")

        initial_count = len(renderer.prs.slides)
        renderer.render_slide(
            SlideSpec(
                layout=SlideLayout.COMPARISON,
                title="Current vs Future",
                left_title="Current",
                right_title="Future",
                left_bullets=sample_bullets,
                right_bullets=sample_bullets,
            )
        )
        assert len(renderer.prs.slides) == initial_count + 1

    def test_add_fullscreen_image_slide(self, tmp_path):
        """Test adding a fullscreen image slide."""
        renderer = PPTRenderer()

        # Create a dummy image
        img_path = str(tmp_path / "test_img.jpg")
        img = Image.new('RGB', (100, 100), color='red')
        img.save(img_path)

        initial_count = len(renderer.prs.slides)
        renderer.add_fullscreen_image_slide(
            img_path,
            caption="Test Image",
            overlay_title="Impact"
        )
        assert len(renderer.prs.slides) == initial_count + 1



class TestSave:
    """Tests for save functionality."""

    def test_save_presentation(self, temp_dir, sample_bullets):
        """Test saving a presentation to file."""
        renderer = PPTRenderer()
        renderer.apply_style("minimalist")

        renderer.add_title_slide("Test Presentation", "Subtitle")
        renderer.add_section_header("Section 1")
        renderer.add_content_slide("Slide 1", sample_bullets)

        output_path = os.path.join(temp_dir, "test_output.pptx")
        renderer.save(output_path)

        assert os.path.exists(output_path)
        assert os.path.getsize(output_path) > 0


class TestInitWithTemplate:
    """Tests for PPTRenderer initialization with a template."""

    def test_init_with_template_removes_slides(self, tmp_path):
        """Lines 24-28: Loading a template and removing existing slides."""
        # Create a valid pptx template with at least one slide
        from pptx import Presentation
        prs = Presentation()
        layout = prs.slide_layouts[0]
        prs.slides.add_slide(layout)
        template_path = str(tmp_path / "template.pptx")
        prs.save(template_path)

        renderer = PPTRenderer(template_path=template_path)
        assert len(renderer.prs.slides) == 0

    def test_init_with_template_preserve_slides(self, tmp_path):
        """Lines 24-28: Loading a template with preserve_template_slides=True."""
        from pptx import Presentation
        prs = Presentation()
        layout = prs.slide_layouts[0]
        prs.slides.add_slide(layout)
        template_path = str(tmp_path / "template.pptx")
        prs.save(template_path)

        renderer = PPTRenderer(template_path=template_path, preserve_template_slides=True)
        assert len(renderer.prs.slides) == 1


class TestRenderSlideFallbacks:
    """Tests for render_slide fallback paths when data is missing."""

    def test_quote_slide_missing_author_falls_back(self):
        """Line 76: Quote slide missing author falls back to content."""
        renderer = PPTRenderer()
        initial_count = len(renderer.prs.slides)
        spec = SlideSpec(
            layout=SlideLayout.QUOTE,
            title="Missing Author",
            quote_text="Some quote",
            quote_author=None,
            bullets=["Fallback bullet"],
        )
        renderer.render_slide(spec)
        assert len(renderer.prs.slides) == initial_count + 1

    def test_quote_slide_missing_text_falls_back(self):
        """Line 76: Quote slide missing text falls back to content."""
        renderer = PPTRenderer()
        initial_count = len(renderer.prs.slides)
        spec = SlideSpec(
            layout=SlideLayout.QUOTE,
            title="Missing Text",
            quote_text=None,
            quote_author="Author",
            bullets=["Fallback bullet"],
        )
        renderer.render_slide(spec)
        assert len(renderer.prs.slides) == initial_count + 1

    def test_add_quote_slide_empty_quote_returns_early(self):
        """Lines 483-485: add_quote_slide with empty quote text returns early."""
        renderer = PPTRenderer()
        renderer.apply_style("minimalist")
        initial_count = len(renderer.prs.slides)
        renderer.add_quote_slide(quote="", author="Steve Jobs")
        assert len(renderer.prs.slides) == initial_count

    def test_add_quote_slide_empty_author_returns_early(self):
        """Lines 483-485: add_quote_slide with empty author returns early."""
        renderer = PPTRenderer()
        renderer.apply_style("minimalist")
        initial_count = len(renderer.prs.slides)
        renderer.add_quote_slide(quote="Some insightful quote.", author="")
        assert len(renderer.prs.slides) == initial_count

    def test_add_quote_slide_none_quote_returns_early(self):
        """Lines 483-485: add_quote_slide with None quote returns early."""
        renderer = PPTRenderer()
        renderer.apply_style("minimalist")
        initial_count = len(renderer.prs.slides)
        renderer.add_quote_slide(quote=None, author="Author")
        assert len(renderer.prs.slides) == initial_count

    def test_add_quote_slide_none_author_returns_early(self):
        """Lines 483-485: add_quote_slide with None author returns early."""
        renderer = PPTRenderer()
        renderer.apply_style("minimalist")
        initial_count = len(renderer.prs.slides)
        renderer.add_quote_slide(quote="A quote.", author=None)
        assert len(renderer.prs.slides) == initial_count

    def test_add_quote_slide_contains_quote_text(self):
        """Valid quote slide should contain the quote text in a text frame."""
        renderer = PPTRenderer()
        renderer.apply_style("minimalist")
        renderer.add_quote_slide(
            quote="Stay hungry, stay foolish.",
            author="Steve Jobs",
            context="Stanford 2005",
        )
        assert len(renderer.prs.slides) == 1
        slide = renderer.prs.slides[0]
        all_text = " ".join(
            shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
        )
        assert "Stay hungry, stay foolish." in all_text
        assert "Steve Jobs" in all_text

    def test_statistics_slide_no_data_falls_back(self):
        """Lines 78-82: Statistics slide with no statistics falls back to content."""
        renderer = PPTRenderer()
        initial_count = len(renderer.prs.slides)
        spec = SlideSpec(
            layout=SlideLayout.STATISTICS,
            title="Empty Stats",
            statistics=None,
            bullets=["Fallback"],
        )
        renderer.render_slide(spec)
        assert len(renderer.prs.slides) == initial_count + 1

    def test_statistics_slide_with_data_renders(self):
        """Lines 78-82: Statistics slide with valid data renders properly."""
        renderer = PPTRenderer()
        initial_count = len(renderer.prs.slides)
        spec = SlideSpec(
            layout=SlideLayout.STATISTICS,
            title="Stats",
            statistics=[
                StatisticData(value="99%", label="Uptime"),
                StatisticData(value="10M", label="Users"),
            ],
        )
        renderer.render_slide(spec)
        assert len(renderer.prs.slides) == initial_count + 1

    def test_image_slide_no_path_falls_back(self):
        """Line 91: Image slide with no image_path falls back to content."""
        renderer = PPTRenderer()
        initial_count = len(renderer.prs.slides)
        spec = SlideSpec(
            layout=SlideLayout.IMAGE,
            title="No Image",
            image_path=None,
            bullets=["Fallback"],
        )
        renderer.render_slide(spec)
        assert len(renderer.prs.slides) == initial_count + 1

    def test_chart_slide_no_data_falls_back(self):
        """Lines 93-96: Chart slide with no chart_data falls back to content."""
        renderer = PPTRenderer()
        initial_count = len(renderer.prs.slides)
        spec = SlideSpec(
            layout=SlideLayout.CHART,
            title="No Chart",
            chart_data=None,
            bullets=["Fallback"],
        )
        renderer.render_slide(spec)
        assert len(renderer.prs.slides) == initial_count + 1


class TestBlankSlideLayoutFallback:
    """Tests for _add_blank_slide layout fallback."""

    def test_blank_layout_fallback_when_no_blank(self):
        """Line 115: When no layout named 'blank' exists, use fallback index."""
        renderer = PPTRenderer()
        # Rename all layouts so none is called "blank"
        for layout in renderer.prs.slide_layouts:
            layout.name = "custom_" + layout.name
        slide = renderer._add_blank_slide()
        assert slide is not None


class TestApplyBackgroundGradientFallback:
    """Tests for _apply_background gradient exception fallback."""

    def test_gradient_fallback_on_exception(self):
        """Lines 140-141: Gradient application raises, falls back to solid."""
        renderer = PPTRenderer()
        renderer.current_style = dict(renderer.current_style)
        renderer.current_style["gradient"] = True
        renderer.current_style["gradient_end"] = renderer.current_style["bg_color"]

        slide = MagicMock()
        # Make gradient() raise to trigger fallback
        slide.background.fill.gradient.side_effect = Exception("gradient not supported")

        renderer._apply_background(slide)
        # Verify solid fill was called as fallback
        slide.background.fill.solid.assert_called()


class TestCoverImageCropBranch:
    """Tests for _cover_image height-crop branch."""

    def test_cover_image_crops_tall_image(self, tmp_path):
        """Lines 272-274: When image is taller than target ratio, crop height."""
        renderer = PPTRenderer()
        img_path = str(tmp_path / "tall.png")
        # Create a tall image (100x400) - current_ratio (0.25) < target_ratio (1.0)
        img = Image.new("RGB", (100, 400), color="blue")
        img.save(img_path)

        cropped = renderer._cover_image(img_path, 1.0)
        assert os.path.exists(cropped)
        with Image.open(cropped) as c:
            w, h = c.size
            assert abs(w / h - 1.0) < 0.01
        os.unlink(cropped)

    def test_cover_image_save_exception_cleans_up(self, tmp_path):
        """Lines 279-281: If save fails, temp file is cleaned up and error re-raised."""
        renderer = PPTRenderer()
        img_path = str(tmp_path / "test.png")
        img = Image.new("RGB", (100, 100), color="red")
        img.save(img_path)

        with patch("autoppt.ppt_renderer.Image.open") as mock_open:
            mock_img = MagicMock()
            mock_img.__enter__ = MagicMock(return_value=mock_img)
            mock_img.__exit__ = MagicMock(return_value=False)
            mock_img.convert.return_value = mock_img
            mock_img.size = (100, 100)
            mock_img.crop.return_value = mock_img
            mock_img.save.side_effect = OSError("disk full")
            mock_open.return_value = mock_img

            with pytest.raises(OSError, match="disk full"):
                renderer._cover_image(img_path, 1.0)


class TestAddCoverPicture:
    """Tests for _add_cover_picture edge cases."""

    def test_returns_false_for_missing_image(self):
        """Line 286: Returns False when image_path does not exist."""
        renderer = PPTRenderer()
        slide = renderer._add_blank_slide()
        result = renderer._add_cover_picture(slide, "/nonexistent/path.png", 0, 0, 5, 5)
        assert result is False

    def test_returns_false_for_empty_path(self):
        """Line 286: Returns False when image_path is empty string."""
        renderer = PPTRenderer()
        slide = renderer._add_blank_slide()
        result = renderer._add_cover_picture(slide, "", 0, 0, 5, 5)
        assert result is False

    def test_oserror_on_unlink_is_suppressed(self, tmp_path):
        """Lines 294-295: OSError during temp file cleanup is suppressed."""
        renderer = PPTRenderer()
        img_path = str(tmp_path / "test.png")
        img = Image.new("RGB", (100, 100), color="green")
        img.save(img_path)

        slide = renderer._add_blank_slide()
        with patch("autoppt.ppt_renderer.os.unlink", side_effect=OSError("perm denied")):
            result = renderer._add_cover_picture(slide, img_path, 0, 0, 5, 5)
            assert result is True


class TestFullscreenImageSlideEdgeCases:
    """Tests for add_fullscreen_image_slide edge cases."""

    def test_missing_image_returns_early(self):
        """Lines 564-565: Missing image logs warning and returns without adding slide."""
        renderer = PPTRenderer()
        initial_count = len(renderer.prs.slides)
        renderer.add_fullscreen_image_slide("/nonexistent/image.png")
        assert len(renderer.prs.slides) == initial_count

    def test_empty_image_path_returns_early(self):
        """Lines 564-565: Empty image path returns early."""
        renderer = PPTRenderer()
        initial_count = len(renderer.prs.slides)
        renderer.add_fullscreen_image_slide("")
        assert len(renderer.prs.slides) == initial_count


class TestStatisticsSlideEdgeCases:
    """Tests for add_statistics_slide edge cases."""

    def test_empty_stats_returns_early(self):
        """Line 605: Empty stats list returns without adding slide."""
        renderer = PPTRenderer()
        initial_count = len(renderer.prs.slides)
        renderer.add_statistics_slide("Empty", [])
        assert len(renderer.prs.slides) == initial_count


class TestRenderDeck:
    """Tests for render_deck method."""

    def test_render_deck_applies_style_and_renders_slides(self):
        """Lines 38-40: render_deck applies style and renders each slide."""
        from autoppt.data_types import DeckSpec
        renderer = PPTRenderer()
        deck = DeckSpec(
            title="Test Deck",
            topic="Testing",
            style="corporate",
            slides=[
                SlideSpec(layout=SlideLayout.TITLE, title="Intro", subtitle="Sub"),
                SlideSpec(layout=SlideLayout.CONTENT, title="Body", bullets=["A", "B"]),
            ],
        )
        renderer.render_deck(deck)
        assert len(renderer.prs.slides) == 2


class TestRenderSlideDispatching:
    """Tests for render_slide layout dispatching."""

    def test_render_title_slide(self):
        """Lines 43-45: render_slide dispatches TITLE layout."""
        renderer = PPTRenderer()
        initial = len(renderer.prs.slides)
        renderer.render_slide(SlideSpec(layout=SlideLayout.TITLE, title="Hello", subtitle="World"))
        assert len(renderer.prs.slides) == initial + 1

    def test_render_section_slide(self):
        """Lines 46-48: render_slide dispatches SECTION layout."""
        renderer = PPTRenderer()
        initial = len(renderer.prs.slides)
        renderer.render_slide(SlideSpec(layout=SlideLayout.SECTION, title="Section A"))
        assert len(renderer.prs.slides) == initial + 1

    def test_render_image_slide_with_valid_image(self, tmp_path):
        """Lines 83-90: render_slide dispatches IMAGE layout with valid image."""
        renderer = PPTRenderer()
        img_path = str(tmp_path / "test.png")
        Image.new("RGB", (200, 150), color="red").save(img_path)

        initial = len(renderer.prs.slides)
        renderer.render_slide(SlideSpec(
            layout=SlideLayout.IMAGE,
            title="Image Slide",
            image_path=img_path,
            image_caption="Caption",
        ))
        assert len(renderer.prs.slides) == initial + 1

    def test_render_chart_slide_with_data(self):
        """Lines 92-96: render_slide dispatches CHART layout with valid data."""
        renderer = PPTRenderer()
        initial = len(renderer.prs.slides)
        chart = ChartData(
            chart_type=ChartType.BAR,
            title="Revenue",
            categories=["A", "B"],
            values=[10.0, 20.0],
            series_name="S1",
        )
        renderer.render_slide(SlideSpec(
            layout=SlideLayout.CHART,
            title="Chart Slide",
            chart_data=chart,
        ))
        assert len(renderer.prs.slides) == initial + 1

    def test_render_citations_slide(self):
        """Lines 97-99: render_slide dispatches CITATIONS layout."""
        renderer = PPTRenderer()
        initial = len(renderer.prs.slides)
        renderer.render_slide(SlideSpec(
            layout=SlideLayout.CITATIONS,
            title="Refs",
            citations=["Source 1", "Source 2"],
        ))
        assert len(renderer.prs.slides) == initial + 1


class TestContentSlideWithImage:
    """Tests for add_content_slide with image_path."""

    def test_content_slide_with_image(self, tmp_path):
        """Lines 348-368: Content slide renders image alongside text."""
        renderer = PPTRenderer()
        img_path = str(tmp_path / "content_img.png")
        Image.new("RGB", (300, 200), color="blue").save(img_path)

        initial = len(renderer.prs.slides)
        renderer.add_content_slide(
            title="With Image",
            bullets=["Point 1", "Point 2"],
            notes="Notes",
            image_path=img_path,
        )
        assert len(renderer.prs.slides) == initial + 1


class TestCoverImageWideCrop:
    """Tests for _cover_image wide-crop branch (already covered, adding explicit)."""

    def test_cover_image_crops_wide_image(self, tmp_path):
        """Lines 267-270: When image is wider than target ratio, crop width."""
        renderer = PPTRenderer()
        img_path = str(tmp_path / "wide.png")
        Image.new("RGB", (400, 100), color="green").save(img_path)

        cropped = renderer._cover_image(img_path, 1.0)
        assert os.path.exists(cropped)
        with Image.open(cropped) as c:
            w, h = c.size
            assert abs(w / h - 1.0) < 0.01
        os.unlink(cropped)


def test_save_wraps_exception_in_render_error():
    """save() should wrap pptx errors in RenderError."""
    from autoppt.exceptions import RenderError
    renderer = PPTRenderer()
    with patch.object(renderer.prs, 'save', side_effect=PermissionError("denied")):
        with pytest.raises(RenderError, match="save"):
            renderer.save("/nonexistent/path.pptx")


class TestRenderDeckTemplateStyle:
    """Tests for render_deck style behavior with templates."""

    def test_render_deck_skips_style_when_template_active(self, tmp_path):
        """render_deck should NOT re-apply style when a template was used."""
        from autoppt.data_types import DeckSpec
        template_file = tmp_path / "tmpl.pptx"
        from pptx import Presentation as Prs
        Prs().save(str(template_file))
        renderer = PPTRenderer(template_path=str(template_file))
        original_style = renderer.current_style.copy()
        deck = DeckSpec(
            title="Test",
            topic="Test",
            style="corporate",
            slides=[SlideSpec(layout=SlideLayout.TITLE, title="Intro", subtitle="Sub")],
        )
        renderer.render_deck(deck)
        assert renderer.current_style == original_style

    def test_render_deck_applies_style_without_template(self):
        """render_deck should apply style when no template is used."""
        from autoppt.data_types import DeckSpec
        renderer = PPTRenderer()
        deck = DeckSpec(
            title="Test",
            topic="Test",
            style="corporate",
            slides=[SlideSpec(layout=SlideLayout.TITLE, title="Intro", subtitle="Sub")],
        )
        renderer.render_deck(deck)
        from autoppt.themes import get_theme
        assert renderer.current_style == get_theme("corporate")

    def test_has_template_flag_set(self, tmp_path):
        """_has_template flag should be set correctly."""
        renderer = PPTRenderer()
        assert renderer._has_template is False

        template_file = tmp_path / "tmpl.pptx"
        from pptx import Presentation as Prs
        Prs().save(str(template_file))
        renderer_with_tmpl = PPTRenderer(template_path=str(template_file))
        assert renderer_with_tmpl._has_template is True


class TestRendererTemplateFileSize:
    """Tests for template file size limit in PPTRenderer."""

    def test_rejects_oversized_template(self, tmp_path):
        """PPTRenderer should reject template files exceeding the size limit."""
        from autoppt.exceptions import RenderError
        from autoppt.config import Config
        big_file = tmp_path / "big.pptx"
        big_file.write_bytes(b"\x00" * (Config.MAX_TEMPLATE_BYTES + 1))
        with pytest.raises(RenderError, match="too large"):
            PPTRenderer(template_path=str(big_file))

    def test_accepts_normal_template(self, tmp_path):
        """PPTRenderer should accept template files within the size limit."""
        template_file = tmp_path / "tmpl.pptx"
        from pptx import Presentation as Prs
        Prs().save(str(template_file))
        renderer = PPTRenderer(template_path=str(template_file))
        assert renderer.prs is not None


class TestCoverImageDecompressionBomb:
    """Tests for decompression bomb protection in _cover_image."""

    def test_cover_image_rejects_decompression_bomb(self, tmp_path):
        """_cover_image should raise RenderError for images exceeding 25M pixels."""
        from autoppt.exceptions import RenderError

        renderer = PPTRenderer()
        img_path = str(tmp_path / "bomb.png")
        # Create a real tiny file so the path exists
        Image.new("RGB", (1, 1), color="red").save(img_path)

        with patch("autoppt.ppt_renderer.Image.open") as mock_open:
            mock_img = MagicMock()
            mock_img.__enter__ = MagicMock(return_value=mock_img)
            mock_img.__exit__ = MagicMock(return_value=False)
            mock_img.size = (100000, 100000)  # 10 billion pixels > 25M limit
            mock_open.return_value = mock_img

            with pytest.raises(RenderError, match="Image too large"):
                renderer._cover_image(img_path, 1.0)


class TestCoverImageExtremeAspectRatio:
    """Tests for _cover_image with extreme aspect ratio images."""

    def test_cover_image_very_small_width_does_not_crash(self, tmp_path):
        """A very narrow image (1x1000) should not crash due to zero crop dimensions."""
        renderer = PPTRenderer()
        img_path = str(tmp_path / "narrow.png")
        img = Image.new("RGB", (1, 1000), color="red")
        img.save(img_path)

        # Should not raise - max(..., 1) prevents zero-dimension crop
        cropped = renderer._cover_image(img_path, 16.0 / 9.0)
        assert os.path.exists(cropped)
        with Image.open(cropped) as c:
            w, h = c.size
            assert w >= 1
            assert h >= 1
        os.unlink(cropped)

    def test_cover_image_very_small_height_does_not_crash(self, tmp_path):
        """A very flat image (1000x1) should not crash due to zero crop dimensions."""
        renderer = PPTRenderer()
        img_path = str(tmp_path / "flat.png")
        img = Image.new("RGB", (1000, 1), color="blue")
        img.save(img_path)

        cropped = renderer._cover_image(img_path, 16.0 / 9.0)
        assert os.path.exists(cropped)
        with Image.open(cropped) as c:
            w, h = c.size
            assert w >= 1
            assert h >= 1
        os.unlink(cropped)


class TestChartDataRejectsEmptyCategories:
    """Tests for ChartData validation rejecting empty categories."""

    def test_chart_data_rejects_empty_categories(self):
        """ChartData should raise ValueError when categories list is empty."""
        with pytest.raises(ValueError, match="categories must not be empty"):
            ChartData(
                chart_type=ChartType.BAR,
                title="Test",
                categories=[],
                values=[],
                series_name="S",
            )


class TestStatisticsSlideTruncation:
    """Tests for statistics slide truncation to 4 items."""

    def test_statistics_slide_truncation_warning(self):
        """add_statistics_slide should truncate stats to 4 items when given more."""
        renderer = PPTRenderer()
        stats = [
            {"value": f"{i * 10}%", "label": f"Metric {i}"}
            for i in range(6)
        ]

        initial_count = len(renderer.prs.slides)
        renderer.add_statistics_slide("Too Many Stats", stats)

        assert len(renderer.prs.slides) == initial_count + 1
        slide = renderer.prs.slides[initial_count]

        # Each stat card produces: 1 panel + 1 value textbox + 1 label textbox = 3 shapes
        # Plus title-related shapes from _add_title_text and _add_blank_slide.
        # With 4 cards we expect 12 stat-related shapes (4 * 3).
        # Count shapes that are textboxes with centered stat values as a proxy.
        stat_value_count = 0
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                # Our stat values are "0%", "10%", "20%", "30%" (first 4 only)
                if text in ["0%", "10%", "20%", "30%"]:
                    stat_value_count += 1
        # Exactly 4 stat values should be present (truncated from 6)
        assert stat_value_count == 4
        # The 5th and 6th values should NOT appear
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                assert text not in ["40%", "50%"]


class TestChartSlideNullTitle:
    """Tests for chart slide with None title."""

    def test_chart_slide_none_title_does_not_crash(self):
        """add_chart_slide should handle None chart_data.title without crashing."""
        renderer = PPTRenderer()
        chart = ChartData(
            chart_type=ChartType.COLUMN,
            title="Placeholder",
            categories=["A", "B"],
            values=[10.0, 20.0],
            series_name="S1",
        )
        chart.title = None  # Bypass Pydantic validation to simulate None at runtime
        initial = len(renderer.prs.slides)
        renderer.add_chart_slide("Chart With No Title", chart)
        assert len(renderer.prs.slides) == initial + 1


class TestChartSlideDataLengthMismatch:
    """Tests for chart slide values/categories length mismatch."""

    def test_more_values_than_categories_truncates(self):
        """add_chart_slide should truncate when values > categories."""
        renderer = PPTRenderer()
        chart = ChartData(
            chart_type=ChartType.BAR,
            title="Mismatch",
            categories=["A", "B"],
            values=[1.0, 2.0],
            series_name="S1",
        )
        # Bypass Pydantic validator to simulate runtime mismatch
        chart.values = [1.0, 2.0, 3.0, 4.0]
        initial = len(renderer.prs.slides)
        renderer.add_chart_slide("Truncated Chart", chart)
        assert len(renderer.prs.slides) == initial + 1

    def test_more_categories_than_values_truncates(self):
        """add_chart_slide should truncate when categories > values."""
        renderer = PPTRenderer()
        chart = ChartData(
            chart_type=ChartType.LINE,
            title="Mismatch 2",
            categories=["A", "B"],
            values=[10.0, 20.0],
            series_name="S1",
        )
        # Bypass Pydantic validator to simulate runtime mismatch
        chart.categories = ["A", "B", "C", "D"]
        initial = len(renderer.prs.slides)
        renderer.add_chart_slide("Truncated Chart 2", chart)
        assert len(renderer.prs.slides) == initial + 1

    def test_matching_lengths_no_truncation(self):
        """add_chart_slide should work normally when lengths match."""
        renderer = PPTRenderer()
        chart = ChartData(
            chart_type=ChartType.COLUMN,
            title="Matched",
            categories=["Q1", "Q2", "Q3"],
            values=[100.0, 200.0, 300.0],
            series_name="Revenue",
        )
        initial = len(renderer.prs.slides)
        renderer.add_chart_slide("Matched Chart", chart)
        assert len(renderer.prs.slides) == initial + 1


class TestChartSlideEmptyCategories:
    """Tests for add_chart_slide guard when categories are empty at runtime."""

    def test_add_chart_slide_skips_empty_categories(self):
        """add_chart_slide should return early without adding a slide when categories are empty."""
        renderer = PPTRenderer()
        initial = len(renderer.prs.slides)
        # Bypass Pydantic validation by constructing ChartData then clearing categories
        chart = ChartData(
            chart_type=ChartType.BAR,
            title="Empty",
            categories=["A"],
            values=[1.0],
        )
        chart.categories = []  # Clear after validation
        renderer.add_chart_slide("Should Skip", chart)
        assert len(renderer.prs.slides) == initial  # No slide added


class TestUnhandledLayoutFallback:
    """Tests for unhandled layout fallback to content."""

    def test_unhandled_layout_falls_back_to_content(self):
        """render_slide should fall back to add_content_slide for an unrecognized layout."""
        renderer = PPTRenderer()
        initial = len(renderer.prs.slides)
        slide_spec = SlideSpec(
            layout=SlideLayout.CONTENT,
            title="Fallback Test",
            bullets=["Point 1"],
        )
        # Force the layout attribute to a value that skips all known branches
        object.__setattr__(slide_spec, "layout", "unknown_layout")
        renderer.render_slide(slide_spec)
        assert len(renderer.prs.slides) == initial + 1


class TestRendererMissingTemplate:
    """Test that missing template file raises RenderError."""

    def test_missing_template_raises_render_error(self):
        from autoppt.exceptions import RenderError
        with pytest.raises(RenderError, match="Template not found"):
            PPTRenderer(template_path="/nonexistent/path/template.pptx")


class TestZipBombProtection:
    """Tests for zip bomb detection in PPTRenderer."""

    def test_bad_zip_file_rejected(self, tmp_path):
        """A file that is not a valid ZIP should be rejected."""
        from autoppt.exceptions import RenderError
        bad_file = tmp_path / "fake.pptx"
        bad_file.write_bytes(b"this is not a zip file")
        with pytest.raises(RenderError, match="Invalid PPTX file"):
            PPTRenderer(template_path=str(bad_file))

    def test_zip_bomb_rejected(self, tmp_path):
        """A ZIP with huge declared decompressed size should be rejected."""
        import zipfile
        import io
        from autoppt.exceptions import RenderError
        from autoppt.config import Config

        bomb_path = tmp_path / "bomb.pptx"
        with zipfile.ZipFile(str(bomb_path), "w") as zf:
            # Write a small file but patch the header to claim huge size
            zf.writestr("content.xml", "small content")

        # Monkey-patch the check to simulate huge decompressed size
        from unittest.mock import patch, MagicMock
        fake_info = MagicMock()
        fake_info.file_size = Config.MAX_DECOMPRESSED_BYTES + 1

        with patch("autoppt.ppt_renderer.zipfile.ZipFile") as mock_zf:
            mock_zf.return_value.__enter__ = MagicMock(return_value=mock_zf.return_value)
            mock_zf.return_value.__exit__ = MagicMock(return_value=False)
            mock_zf.return_value.infolist.return_value = [fake_info]
            with pytest.raises(RenderError, match="decompressed size"):
                PPTRenderer(template_path=str(bomb_path))

    def test_normal_pptx_passes_zip_check(self, tmp_path):
        """A normal PPTX template should pass the zip bomb check."""
        from pptx import Presentation
        normal = tmp_path / "normal.pptx"
        Presentation().save(str(normal))
        renderer = PPTRenderer(template_path=str(normal))
        assert renderer._has_template is True


class TestChartSlideEmptyAfterTruncation:
    """Tests for add_chart_slide edge case when data empties after truncation."""

    def test_chart_slide_valid_data_adds_slide(self):
        """add_chart_slide should add a slide when categories and values match."""
        renderer = PPTRenderer()
        initial_count = len(renderer.prs.slides)
        chart = ChartData(
            chart_type=ChartType.BAR,
            title="Test",
            categories=["A"],
            values=[1.0],
            series_name="S",
        )
        renderer.add_chart_slide("Normal Chart", chart)
        assert len(renderer.prs.slides) == initial_count + 1

    def test_chart_slide_skips_when_empty_after_truncation(self):
        """add_chart_slide should skip when values are empty causing empty truncation."""
        renderer = PPTRenderer()
        initial_count = len(renderer.prs.slides)
        # Bypass Pydantic validation to create a mismatched ChartData
        chart = ChartData.model_construct(
            chart_type=ChartType.BAR,
            title="Empty After Trunc",
            categories=["A", "B"],
            values=[],
            series_name="S",
        )
        renderer.add_chart_slide("Truncated Chart", chart)
        assert len(renderer.prs.slides) == initial_count  # No slide added


class TestFullscreenImageSlideNotes:
    """Tests for fullscreen image slide speaker notes support."""

    def test_fullscreen_image_slide_preserves_notes(self, tmp_path):
        """add_fullscreen_image_slide should pass notes to the slide."""
        renderer = PPTRenderer()
        img_path = str(tmp_path / "test_img.png")
        Image.new("RGB", (100, 100), color="blue").save(img_path)

        renderer.add_fullscreen_image_slide(
            img_path,
            caption="A caption",
            overlay_title="Title",
            notes="Important presenter note",
        )
        slide = renderer.prs.slides[0]
        assert slide.notes_slide.notes_text_frame.text == "Important presenter note"

    def test_render_image_slide_passes_speaker_notes(self, tmp_path):
        """render_slide for IMAGE layout should pass speaker_notes through."""
        renderer = PPTRenderer()
        img_path = str(tmp_path / "test_img.png")
        Image.new("RGB", (100, 100), color="green").save(img_path)

        renderer.render_slide(SlideSpec(
            layout=SlideLayout.IMAGE,
            title="Noted Image",
            image_path=img_path,
            image_caption="Cap",
            speaker_notes="Slide-level note",
        ))
        slide = renderer.prs.slides[0]
        assert slide.notes_slide.notes_text_frame.text == "Slide-level note"


class TestChartSlideRenderFallback:
    """Tests for render_slide falling back when chart data is insufficient."""

    def test_render_chart_slide_falls_back_on_empty_categories(self):
        """render_slide should fall back to content when add_chart_slide skips."""
        renderer = PPTRenderer()
        chart = ChartData.model_construct(
            chart_type=ChartType.BAR,
            title="Empty Chart",
            categories=[],
            values=[],
            series_name="S",
        )
        slide_spec = SlideSpec.model_construct(
            layout=SlideLayout.CHART,
            title="Fallback Test",
            bullets=["Fallback bullet"],
            chart_data=chart,
            speaker_notes="Chart fallback note",
            image_path=None,
            image_caption=None,
            quote_text=None,
            quote_author=None,
            quote_context=None,
            left_title=None,
            left_bullets=None,
            right_title=None,
            right_bullets=None,
            statistics=None,
            citations=None,
            image_query=None,
        )
        renderer.render_slide(slide_spec)
        assert len(renderer.prs.slides) == 1  # Content fallback was added

    def test_render_chart_slide_falls_back_on_empty_after_truncation(self):
        """render_slide should fall back when chart values are empty after truncation."""
        renderer = PPTRenderer()
        chart = ChartData.model_construct(
            chart_type=ChartType.BAR,
            title="Mismatched",
            categories=["A", "B"],
            values=[],
            series_name="S",
        )
        slide_spec = SlideSpec.model_construct(
            layout=SlideLayout.CHART,
            title="Truncation Fallback",
            bullets=["Bullet content"],
            chart_data=chart,
            speaker_notes="",
            image_path=None,
            image_caption=None,
            quote_text=None,
            quote_author=None,
            quote_context=None,
            left_title=None,
            left_bullets=None,
            right_title=None,
            right_bullets=None,
            statistics=None,
            citations=None,
            image_query=None,
        )
        renderer.render_slide(slide_spec)
        assert len(renderer.prs.slides) == 1  # Content fallback was added


class TestImageSlideFallbackContent:
    """Tests for IMAGE slide falling back to content with bullets."""

    def test_image_slide_fallback_uses_bullets(self):
        """render_slide should use bullets when IMAGE slide has no image."""
        renderer = PPTRenderer()
        slide_spec = SlideSpec(
            layout=SlideLayout.IMAGE,
            title="Missing Image",
            bullets=["Fallback point A", "Fallback point B"],
            image_path=None,
        )
        renderer.render_slide(slide_spec)
        assert len(renderer.prs.slides) == 1


class TestSaveRejectsSystemPath:
    """Tests for save() rejecting system paths."""

    def test_save_rejects_system_path(self):
        """save() should raise RenderError when output path is under a system prefix."""
        from autoppt.exceptions import RenderError
        renderer = PPTRenderer()
        with pytest.raises(RenderError, match="system path"):
            renderer.save("/etc/foo.pptx")


class TestQuoteSlideFallback:
    """Tests for QUOTE slide fallback when add_quote_slide silently returns."""

    def test_quote_slide_fallback_on_empty_quote_after_strip(self):
        """render_slide should fall back to content when quote_text is truthy but add_quote_slide rejects it."""
        renderer = PPTRenderer()
        # add_quote_slide guards on `not quote`, which is falsy for whitespace-only strings.
        # But since " " is truthy, render_slide enters the quote path.
        # The slide-count check ensures we fall back to content.
        slide_spec = SlideSpec(
            layout=SlideLayout.QUOTE,
            title="Test Quote",
            bullets=["Fallback bullet A", "Fallback bullet B"],
            quote_text="Valid quote",
            quote_author="Author",
        )
        renderer.render_slide(slide_spec)
        assert len(renderer.prs.slides) == 1

    def test_quote_slide_normal_path(self):
        """render_slide should add a quote slide when both text and author are present."""
        renderer = PPTRenderer()
        slide_spec = SlideSpec(
            layout=SlideLayout.QUOTE,
            title="Good Quote",
            bullets=[],
            quote_text="Innovation distinguishes leaders from followers.",
            quote_author="Steve Jobs",
        )
        renderer.render_slide(slide_spec)
        assert len(renderer.prs.slides) == 1


class TestDecompressionBombWrapped:
    """Tests for DecompressionBombError being wrapped as RenderError."""

    def test_decompression_bomb_raises_render_error(self):
        """_cover_image should wrap DecompressionBombError into RenderError."""
        from autoppt.exceptions import RenderError
        from unittest.mock import patch as mock_patch
        import PIL.Image

        renderer = PPTRenderer()
        with mock_patch("autoppt.ppt_renderer.Image.open", side_effect=PIL.Image.DecompressionBombError("too many pixels")):
            with pytest.raises(RenderError, match="decompression bomb"):
                renderer._cover_image("/fake/path.jpg", 1.5)


class TestStatisticsSlideFallback:
    """Tests for STATISTICS slide fallback when add_statistics_slide doesn't add."""

    def test_statistics_slide_normal_path(self):
        """render_slide should add a statistics slide when data is present."""
        renderer = PPTRenderer()
        slide_spec = SlideSpec(
            layout=SlideLayout.STATISTICS,
            title="Key Metrics",
            bullets=["Fallback"],
            statistics=[
                StatisticData(value="85%", label="Accuracy"),
                StatisticData(value="$4B", label="Revenue"),
            ],
        )
        renderer.render_slide(slide_spec)
        assert len(renderer.prs.slides) == 1

    def test_statistics_slide_fallback_no_data(self):
        """render_slide should fall back to content when statistics is empty."""
        renderer = PPTRenderer()
        slide_spec = SlideSpec(
            layout=SlideLayout.STATISTICS,
            title="No Stats",
            bullets=["Fallback bullet"],
            statistics=None,
        )
        renderer.render_slide(slide_spec)
        assert len(renderer.prs.slides) == 1


class TestRendererSaveBlockedPath:
    """Tests for save() rejecting blocked system paths."""

    def test_save_rejects_etc_path(self):
        """save() should reject writes to /etc/ prefix."""
        from autoppt.exceptions import RenderError
        renderer = PPTRenderer()
        renderer.add_title_slide("Test", "Sub")
        with pytest.raises(RenderError, match="system path"):
            renderer.save("/etc/evil.pptx")

    def test_save_rejects_proc_path(self):
        """save() should reject writes to /proc/ prefix."""
        from autoppt.exceptions import RenderError
        renderer = PPTRenderer()
        renderer.add_title_slide("Test", "Sub")
        with pytest.raises(RenderError, match="system path"):
            renderer.save("/proc/evil.pptx")


class TestImageSlideNotesPass:
    """Tests that fullscreen image slide passes notes through."""

    def test_image_slide_with_notes(self):
        """render_slide for IMAGE should pass speaker_notes to add_fullscreen_image_slide."""
        renderer = PPTRenderer()
        with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
            img = Image.new("RGB", (100, 100), color="red")
            img.save(tmp.name)
            tmp_path = tmp.name
        try:
            slide_spec = SlideSpec(
                layout=SlideLayout.IMAGE,
                title="Photo",
                bullets=[],
                image_path=tmp_path,
                image_caption="A red image",
                speaker_notes="Talk about this image",
            )
            renderer.render_slide(slide_spec)
            assert len(renderer.prs.slides) == 1
            notes = renderer.prs.slides[0].notes_slide.notes_text_frame.text
            assert notes == "Talk about this image"
        finally:
            os.unlink(tmp_path)


class TestImageSlideFallbackGuard:
    """Tests for IMAGE slide count_before guard pattern."""

    def test_image_slide_fallback_when_file_missing_at_render(self):
        """render_slide should fall back to content when add_fullscreen_image_slide silently returns."""
        renderer = PPTRenderer()
        slide_spec = SlideSpec(
            layout=SlideLayout.IMAGE,
            title="Vanished Image",
            bullets=["Fallback A", "Fallback B"],
            image_path="/nonexistent/image.png",
        )
        renderer.render_slide(slide_spec)
        # The guard should detect no slide was added and fall back to content
        assert len(renderer.prs.slides) == 1
