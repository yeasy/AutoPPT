import tempfile
from pathlib import Path
from unittest.mock import patch, MagicMock

import pytest
from pptx import Presentation
from PIL import Image, ImageFont

from autoppt.data_types import DeckSpec, SlideLayout, SlideSpec
from autoppt.sample_library import (
    SampleDefinition,
    _get_card_lines,
    _gradient_image,
    _load_font,
    _real_preview_image_for_sample,
    _draw_showcase_card,
    build_sample_deck,
    get_sample_definition,
    get_sample_definitions,
    render_readme_showcase_previews,
    render_sample,
)


def test_sample_definitions_and_builders_in_sync():
    """SAMPLE_DEFINITIONS and _SAMPLE_BUILDERS must have matching IDs."""
    from autoppt.sample_library import SAMPLE_DEFINITIONS, _SAMPLE_BUILDERS
    defined_ids = {d.sample_id for d in SAMPLE_DEFINITIONS}
    builder_ids = set(_SAMPLE_BUILDERS.keys())
    assert defined_ids == builder_ids, (
        f"defined only: {defined_ids - builder_ids}, builders only: {builder_ids - defined_ids}"
    )


def test_sample_library_builds_all_defined_decks():
    definitions = get_sample_definitions()

    assert len(definitions) >= 11
    for definition in definitions:
        with tempfile.TemporaryDirectory(prefix="test-sample-") as asset_dir:
            deck = build_sample_deck(definition.sample_id, asset_dir=asset_dir)
            assert deck.title
            assert deck.style == definition.style
            assert deck.language == definition.language
            assert len(deck.slides) >= 6


def test_sample_library_renders_sample_pptx(tmp_path):
    output_path = render_sample("en_visual_showcase", tmp_path)

    assert output_path.exists()
    presentation = Presentation(str(output_path))
    assert len(presentation.slides) >= 6


def test_sample_library_renders_readme_previews(tmp_path):
    outputs = render_readme_showcase_previews(tmp_path)

    assert len(outputs) == 2
    for output in outputs:
        assert output.exists()
        assert output.stat().st_size > 0


def test_sample_library_renders_readme_previews_with_real_preview_path(tmp_path, monkeypatch):
    preview_source = tmp_path / "slide-1.jpg"
    Image.new("RGB", (1200, 675), color=(24, 42, 88)).save(preview_source)

    def fake_check_dependencies():
        return True, []

    def fake_render_sample(sample_id, output_dir):
        output = output_dir / f"{sample_id}.pptx"
        output.write_text("fake")
        return output

    def fake_convert_to_pdf(pptx_path, output_dir):
        pdf_path = output_dir / f"{pptx_path.stem}.pdf"
        pdf_path.write_text("fake")
        return pdf_path

    def fake_convert_pdf_to_images(pdf_path, output_dir):
        return [preview_source]

    monkeypatch.setattr("autoppt.sample_library.check_dependencies", fake_check_dependencies)
    monkeypatch.setattr("autoppt.sample_library.render_sample", fake_render_sample)
    monkeypatch.setattr("autoppt.sample_library.convert_to_pdf", fake_convert_to_pdf)
    monkeypatch.setattr("autoppt.sample_library.convert_pdf_to_images", fake_convert_pdf_to_images)

    outputs = render_readme_showcase_previews(tmp_path)
    assert len(outputs) == 2
    assert all(path.exists() for path in outputs)


# --- Tests for get_sample_definitions with category filter (line 557) ---


def test_get_sample_definitions_filters_by_category():
    showcase_defs = get_sample_definitions(category="showcase")
    assert len(showcase_defs) >= 1
    assert all(d.category == "showcase" for d in showcase_defs)

    feature_defs = get_sample_definitions(category="feature")
    assert len(feature_defs) >= 1
    assert all(d.category == "feature" for d in feature_defs)

    # Filtering should return fewer results than "all"
    all_defs = get_sample_definitions(category="all")
    assert len(showcase_defs) < len(all_defs)


def test_get_sample_definitions_unknown_category_returns_empty():
    result = get_sample_definitions(category="nonexistent_category_xyz")
    assert result == []


# --- Tests for get_sample_definition KeyError (line 564) ---


def test_get_sample_definition_raises_key_error_for_unknown_id():
    with pytest.raises(KeyError, match="Unknown sample id"):
        get_sample_definition("totally_bogus_sample_id")


def test_get_sample_definition_returns_known_id():
    defn = get_sample_definition("en_tech")
    assert defn.sample_id == "en_tech"
    assert defn.category == "showcase"


# --- Tests for build_sample_deck ValueError when asset_dir is None (line 578) ---


def test_build_sample_deck_raises_value_error_when_asset_dir_is_none():
    with pytest.raises(ValueError, match="asset_dir is required"):
        build_sample_deck("en_tech", asset_dir=None)


# --- Tests for _get_card_lines fallbacks (lines 282-283, 286) ---


def _make_deck(slides: list[SlideSpec], topic: str = "Test Topic") -> DeckSpec:
    return DeckSpec(title="Test", topic=topic, slides=slides)


def test_get_card_lines_returns_bullets_when_present():
    slide = SlideSpec(layout=SlideLayout.CONTENT, bullets=["a", "b", "c", "d"])
    lines = _get_card_lines(_make_deck([slide]))
    assert lines == ["a", "b", "c"]


def test_get_card_lines_falls_back_to_left_right_bullets():
    slide = SlideSpec(
        layout=SlideLayout.TWO_COLUMN,
        left_bullets=["L1", "L2", "L3"],
        right_bullets=["R1", "R2"],
    )
    lines = _get_card_lines(_make_deck([slide]))
    assert lines == ["L1", "L2", "R1"]


def test_get_card_lines_left_bullets_only():
    slide = SlideSpec(
        layout=SlideLayout.TWO_COLUMN,
        left_bullets=["L1"],
        right_bullets=[],
    )
    lines = _get_card_lines(_make_deck([slide]))
    assert lines == ["L1"]


def test_get_card_lines_right_bullets_only():
    slide = SlideSpec(
        layout=SlideLayout.TWO_COLUMN,
        left_bullets=[],
        right_bullets=["R1", "R2"],
    )
    lines = _get_card_lines(_make_deck([slide]))
    assert lines == ["R1"]


def test_get_card_lines_falls_back_to_quote_text():
    slide = SlideSpec(layout=SlideLayout.QUOTE, quote_text="A famous quote")
    lines = _get_card_lines(_make_deck([slide]))
    assert lines == ["A famous quote"]


def test_get_card_lines_returns_topic_when_no_content():
    slide = SlideSpec(layout=SlideLayout.TITLE, title="Title Only")
    lines = _get_card_lines(_make_deck([slide], topic="My Topic"))
    assert lines == ["My Topic"]


def test_get_card_lines_returns_topic_for_empty_slides():
    lines = _get_card_lines(_make_deck([], topic="Fallback Topic"))
    assert lines == ["Fallback Topic"]


# --- Test for card drawing overflow guard (line 363) ---


def test_draw_showcase_card_overflow_guard():
    """Card drawing should not crash when many bullets exceed card height."""
    many_bullets = [f"Bullet point number {i} with enough text to wrap" for i in range(30)]
    slide = SlideSpec(layout=SlideLayout.CONTENT, bullets=many_bullets)
    deck = _make_deck([slide], topic="Overflow Test")
    deck_with_style = deck.model_copy(update={"style": "minimalist"})

    definition = SampleDefinition(
        sample_id="test",
        filename="test.pptx",
        category="test",
        title="Test Card",
        topic="Test",
        language="English",
        style="minimalist",
        description="Test description",
    )

    # Use a very small card height to trigger the overflow guard quickly
    card = _draw_showcase_card(definition, deck_with_style, card_size=(440, 200), locale="en")
    assert card.size == (440, 200)


# --- Tests for _gradient_image non-positive size (line 173) ---


def test_gradient_image_nonpositive_size():
    with pytest.raises(ValueError, match="Image size must be positive"):
        _gradient_image((0, 100), (0, 0, 0), (255, 255, 255))


def test_gradient_image_nonpositive_height():
    with pytest.raises(ValueError, match="Image size must be positive"):
        _gradient_image((100, -1), (0, 0, 0), (255, 255, 255))


# --- Tests for _load_font fallback to default bitmap font (lines 230-236) ---


def test_load_font_all_fail_uses_default():
    default_font = ImageFont.load_default()
    with patch.object(ImageFont, "truetype", side_effect=OSError("no font")), \
         patch.object(ImageFont, "load_default", return_value=default_font) as mock_default:
        font = _load_font(["/fake/font.ttf", "/other/font.otf"], size=24)
    # Should fall back to the default bitmap font
    mock_default.assert_called_once()
    assert font is default_font


# --- Tests for _real_preview_image_for_sample edge cases (lines 443, 446, 453) ---


def _make_simple_deck() -> DeckSpec:
    return DeckSpec(
        title="Test",
        topic="Test Topic",
        slides=[SlideSpec(layout=SlideLayout.CONTENT, title="Slide 1", bullets=["a"])],
    )


@patch("autoppt.sample_library.check_dependencies", return_value=(True, []))
@patch("autoppt.sample_library.render_sample")
@patch("autoppt.sample_library.convert_to_pdf", return_value=None)
def test_generate_real_preview_convert_to_pdf_returns_none(
    mock_convert_to_pdf, mock_render_sample, mock_check_deps, tmp_path
):
    mock_render_sample.return_value = tmp_path / "test.pptx"
    result = _real_preview_image_for_sample("test_id", _make_simple_deck(), cache_dir=tmp_path)
    assert result is None


@patch("autoppt.sample_library.check_dependencies", return_value=(True, []))
@patch("autoppt.sample_library.render_sample")
@patch("autoppt.sample_library.convert_to_pdf")
@patch("autoppt.sample_library.convert_pdf_to_images", return_value=[])
def test_generate_real_preview_empty_images(
    mock_images, mock_convert_pdf, mock_render_sample, mock_check_deps, tmp_path
):
    mock_render_sample.return_value = tmp_path / "test.pptx"
    mock_convert_pdf.return_value = tmp_path / "test.pdf"
    result = _real_preview_image_for_sample("test_id", _make_simple_deck(), cache_dir=tmp_path)
    assert result is None


@patch("autoppt.sample_library.check_dependencies", return_value=(True, []))
@patch("autoppt.sample_library.render_sample")
@patch("autoppt.sample_library.convert_to_pdf")
@patch("autoppt.sample_library.convert_pdf_to_images")
def test_generate_real_preview_writes_to_cache_dir(
    mock_images, mock_convert_pdf, mock_render_sample, mock_check_deps, tmp_path
):
    mock_render_sample.return_value = tmp_path / "test.pptx"
    pdf_path = tmp_path / "test.pdf"
    pdf_path.write_text("fake")
    mock_convert_pdf.return_value = pdf_path

    preview_img = tmp_path / "slide-1.jpg"
    Image.new("RGB", (800, 600), color=(100, 100, 100)).save(preview_img)
    mock_images.return_value = [preview_img]

    cache = tmp_path / "cache"
    result = _real_preview_image_for_sample("test_id", _make_simple_deck(), cache_dir=cache)
    assert result is not None
    assert result.parent == cache
