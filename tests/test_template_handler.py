"""
Unit tests for TemplateHandler.
"""
import pytest
import os
from pptx import Presentation
from autoppt.template_handler import TemplateHandler

class TestTemplateHandler:
    """Tests for TemplateHandler."""

    @pytest.fixture
    def sample_template(self, tmp_path):
        """Create a temporary real PPTX file."""
        path = tmp_path / "template.pptx"
        prs = Presentation()
        # Add slide master layouts (default has 11 layouts)
        prs.save(path)
        return path

    def test_init_valid_template(self, sample_template):
        """Test initialization with valid template."""
        handler = TemplateHandler(str(sample_template))
        assert handler.prs is not None
        assert isinstance(handler.layouts, dict)

    def test_init_invalid_path(self):
        """Test initialization with invalid path."""
        with pytest.raises(FileNotFoundError):
            TemplateHandler("nonexistent.pptx")

    def test_init_rejects_path_traversal(self):
        """Path traversal in template path must be rejected."""
        with pytest.raises(ValueError, match="Path traversal"):
            TemplateHandler("../../etc/passwd.pptx")

    def test_init_rejects_system_path(self):
        """System paths must be rejected even without traversal."""
        with pytest.raises(ValueError, match="system path"):
            TemplateHandler("/etc/pptx/template.pptx")

    def test_init_rejects_sensitive_path_ssh(self):
        """Paths containing .ssh/ must be rejected."""
        with pytest.raises(ValueError, match="sensitive path"):
            TemplateHandler("/home/user/.ssh/template.pptx")

    def test_init_rejects_sensitive_path_docker(self):
        """Paths containing .docker/ must be rejected."""
        with pytest.raises(ValueError, match="sensitive path"):
            TemplateHandler("/home/user/.docker/template.pptx")

    def test_analyze_layouts(self, sample_template):
        """Test layout analysis."""
        handler = TemplateHandler(str(sample_template))
        layouts = handler.layouts

        # Standard python-pptx presentation has layouts
        # usually 11 layouts in default template
        assert len(layouts) > 0

        # Check structure of layout info
        first_layout = layouts[0]
        assert "name" in first_layout
        assert "placeholders" in first_layout

    def test_get_layout_by_name(self, sample_template):
        """Test finding layout by name."""
        handler = TemplateHandler(str(sample_template))

        # Default layouts have names like "Title Slide", "Title and Content"
        # We need to check what names python-pptx default template uses
        # Usually: "Title Slide", "Title and Content", "Section Header", etc.

        # The default blank presentation might have different names or empty names depending on python-pptx version
        # But generally they match standard Office names.

        found_idx = handler.get_layout_by_name("Title")
        # Default python-pptx templates always have a layout containing "Title"
        assert found_idx is not None
        assert isinstance(found_idx, int)

    def test_get_best_layout_for_type(self, sample_template):
        """Test getting best layout for type."""
        handler = TemplateHandler(str(sample_template))

        idx = handler.get_best_layout_for_type("title")
        assert isinstance(idx, int)

        idx_content = handler.get_best_layout_for_type("content")
        assert isinstance(idx_content, int)

    def test_extract_text_content_no_markitdown(self, sample_template):
        """Test extract_text_content returns empty when markitdown unavailable."""
        from unittest.mock import patch
        handler = TemplateHandler(str(sample_template))
        with patch("autoppt.template_handler.HAS_MARKITDOWN", False):
            content = handler.extract_text_content()
        assert content == ""

    def test_extract_text_content_with_markitdown(self, sample_template):
        """Test extract_text_content with mocked markitdown."""
        from unittest.mock import patch, MagicMock
        import autoppt.template_handler as th_mod
        handler = TemplateHandler(str(sample_template))
        mock_md_cls = MagicMock()
        mock_result = MagicMock()
        mock_result.text_content = "Extracted content"
        mock_md_cls.return_value.convert.return_value = mock_result
        with patch.object(th_mod, "HAS_MARKITDOWN", True), \
             patch.object(th_mod, "MarkItDown", mock_md_cls, create=True):
            content = handler.extract_text_content()
        assert content == "Extracted content"

    def test_extract_text_content_error(self, sample_template):
        """Test extract_text_content handles conversion errors gracefully."""
        from unittest.mock import patch, MagicMock
        import autoppt.template_handler as th_mod
        handler = TemplateHandler(str(sample_template))
        mock_md_cls = MagicMock()
        mock_md_cls.return_value.convert.side_effect = RuntimeError("fail")
        with patch.object(th_mod, "HAS_MARKITDOWN", True), \
             patch.object(th_mod, "MarkItDown", mock_md_cls, create=True):
            content = handler.extract_text_content()
        assert content == ""

    def test_get_best_layout_for_type_content_fallback(self, sample_template):
        """Test content type falls back to layout 1 when no name match."""
        handler = TemplateHandler(str(sample_template))
        handler.layouts = {0: {"name": "X", "placeholders": []}, 1: {"name": "Y", "placeholders": []}}
        idx = handler.get_best_layout_for_type("content")
        assert idx == 1

    def test_get_best_layout_for_type_unknown_returns_content_fallback(self, sample_template):
        """Test unknown type returns layout 1 (content) instead of 0 (title)."""
        handler = TemplateHandler(str(sample_template))
        idx = handler.get_best_layout_for_type("nonexistent_type")
        assert idx == 1

    def test_get_layout_by_name_returns_none_for_no_match(self, sample_template):
        """Test get_layout_by_name returns None for no match."""
        handler = TemplateHandler(str(sample_template))
        idx = handler.get_layout_by_name("ZZZZZZZ_NONEXISTENT")
        assert idx is None


def test_has_markitdown_false_when_import_fails():
    """Test that HAS_MARKITDOWN is False when markitdown is not installed.

    This covers lines 12-13 (the except ImportError branch) by reloading
    the module with markitdown blocked from importing.
    """
    import importlib
    import sys
    from unittest.mock import patch

    # Save original module state
    original_module = sys.modules.get("autoppt.template_handler")
    original_markitdown = sys.modules.get("markitdown")

    try:
        # Remove cached module so it gets re-imported
        sys.modules.pop("autoppt.template_handler", None)
        # Block markitdown from importing
        sys.modules["markitdown"] = None

        # Re-import the module; the try/except ImportError path will execute
        import autoppt.template_handler as th

        assert th.HAS_MARKITDOWN is False
    finally:
        # Restore original state
        sys.modules.pop("autoppt.template_handler", None)
        if original_module is not None:
            sys.modules["autoppt.template_handler"] = original_module
        if original_markitdown is not None:
            sys.modules["markitdown"] = original_markitdown
        else:
            sys.modules.pop("markitdown", None)


def test_rejects_non_pptx_file(tmp_path):
    """TemplateHandler should reject non-.pptx files."""
    txt_file = tmp_path / "template.txt"
    txt_file.write_text("not a pptx")
    with pytest.raises(ValueError, match="Invalid template file type"):
        TemplateHandler(str(txt_file))


def test_rejects_docx_file(tmp_path):
    """TemplateHandler should reject .docx files."""
    docx_file = tmp_path / "template.docx"
    docx_file.write_bytes(b"PK\x03\x04fake")
    with pytest.raises(ValueError, match="Invalid template file type"):
        TemplateHandler(str(docx_file))


def test_rejects_oversized_template(tmp_path):
    """TemplateHandler should reject templates exceeding the size limit."""
    from autoppt.config import Config
    big_file = tmp_path / "huge.pptx"
    big_file.write_bytes(b"\x00" * (Config.MAX_TEMPLATE_BYTES + 1))
    with pytest.raises(ValueError, match="too large"):
        TemplateHandler(str(big_file))


class TestTemplateHandlerZipBomb:
    """Tests for zip bomb detection in TemplateHandler."""

    def test_bad_zip_file_rejected(self, tmp_path):
        """A corrupt file should be rejected."""
        bad_file = tmp_path / "corrupt.pptx"
        bad_file.write_bytes(b"not a zip")
        with pytest.raises(ValueError, match="Invalid PPTX file"):
            TemplateHandler(str(bad_file))

    def test_zip_bomb_rejected(self, tmp_path):
        """Zip bomb should be rejected."""
        import zipfile
        from unittest.mock import patch, MagicMock
        from autoppt.config import Config

        bomb_path = tmp_path / "bomb.pptx"
        with zipfile.ZipFile(str(bomb_path), "w") as zf:
            zf.writestr("content.xml", "small")

        fake_info = MagicMock()
        fake_info.file_size = Config.MAX_DECOMPRESSED_BYTES + 1

        with patch("autoppt.template_handler.zipfile.ZipFile") as mock_zf:
            mock_zf.return_value.__enter__ = MagicMock(return_value=mock_zf.return_value)
            mock_zf.return_value.__exit__ = MagicMock(return_value=False)
            mock_zf.return_value.infolist.return_value = [fake_info]
            with pytest.raises(ValueError, match="decompressed size"):
                TemplateHandler(str(bomb_path))


class TestTemplateHandlerBlockedPaths:
    """Tests for BLOCKED_SYSTEM_PREFIXES in TemplateHandler."""

    def test_rejects_etc_path(self):
        """TemplateHandler should reject paths under /etc/."""
        with pytest.raises(ValueError, match="system path"):
            TemplateHandler("/etc/evil.pptx")

    def test_rejects_proc_path(self):
        """TemplateHandler should reject paths under /proc/."""
        with pytest.raises(ValueError, match="system path"):
            TemplateHandler("/proc/self/evil.pptx")

    def test_rejects_dev_path(self):
        """TemplateHandler should reject paths under /dev/."""
        with pytest.raises(ValueError, match="system path"):
            TemplateHandler("/dev/evil.pptx")
