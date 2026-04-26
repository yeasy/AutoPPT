from __future__ import annotations

import logging
import os
import tempfile
import zipfile
from typing import Any

from PIL import Image
Image.MAX_IMAGE_PIXELS = 25_000_000  # Prevent decompression bombs before full decode
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from .config import Config
from .data_types import ChartData, DeckSpec, SlideLayout, SlideSpec
from .exceptions import RenderError
from .themes import get_theme

logger = logging.getLogger(__name__)

def _check_zip_bomb(path: str) -> None:
    """Reject PPTX files whose decompressed content exceeds the safety limit."""
    try:
        with zipfile.ZipFile(path, "r") as zf:
            total = sum(info.file_size for info in zf.infolist())
            if total > Config.MAX_DECOMPRESSED_BYTES:
                raise RenderError(
                    "init",
                    f"Template decompressed size ({total} bytes) exceeds limit ({Config.MAX_DECOMPRESSED_BYTES})",
                )
    except zipfile.BadZipFile as exc:
        raise RenderError("init", f"Invalid PPTX file: {exc}") from exc


class PPTRenderer:
    def __init__(self, template_path: str | None = None, preserve_template_slides: bool = False):
        self._has_template = bool(template_path)
        if template_path:
            if not os.path.isfile(template_path):
                raise RenderError("init", f"Template not found: {template_path}")
            size = os.path.getsize(template_path)
            if size > Config.MAX_TEMPLATE_BYTES:
                raise RenderError("init", f"Template file too large: {size} bytes (max {Config.MAX_TEMPLATE_BYTES})")
            _check_zip_bomb(template_path)
            self.prs = Presentation(template_path)
            if not preserve_template_slides:
                xml_slides = self.prs.slides._sldIdLst
                for slide in list(xml_slides):
                    xml_slides.remove(slide)
        else:
            self.prs = Presentation()
        self.current_style = get_theme("minimalist")

    def apply_style(self, style_name: str) -> None:
        self.current_style = get_theme(style_name)
        logger.info("Applied theme: %s", style_name)

    def render_deck(self, deck_spec: DeckSpec) -> None:
        if not self._has_template:
            self.apply_style(deck_spec.style)
        for slide_spec in deck_spec.slides:
            self.render_slide(slide_spec)

    def render_slide(self, slide_spec: SlideSpec) -> None:
        if slide_spec.layout == SlideLayout.TITLE:
            self.add_title_slide(slide_spec.title, slide_spec.subtitle or "")
            return
        if slide_spec.layout == SlideLayout.SECTION:
            self.add_section_header(slide_spec.title)
            return
        if slide_spec.layout == SlideLayout.TWO_COLUMN:
            self.add_two_column_slide(
                slide_spec.title,
                slide_spec.left_bullets,
                slide_spec.right_bullets,
                left_title=slide_spec.left_title or "",
                right_title=slide_spec.right_title or "",
                notes=slide_spec.speaker_notes or "",
            )
            return
        if slide_spec.layout == SlideLayout.COMPARISON:
            self.add_comparison_slide(
                slide_spec.title,
                {"name": slide_spec.left_title or "Option A", "points": slide_spec.left_bullets},
                {"name": slide_spec.right_title or "Option B", "points": slide_spec.right_bullets},
                notes=slide_spec.speaker_notes or "",
            )
            return
        if slide_spec.layout == SlideLayout.QUOTE:
            if slide_spec.quote_text and slide_spec.quote_author:
                count_before = len(self.prs.slides)
                self.add_quote_slide(
                    slide_spec.quote_text,
                    slide_spec.quote_author,
                    context=slide_spec.quote_context or "",
                    notes=slide_spec.speaker_notes or "",
                )
                if len(self.prs.slides) > count_before:
                    return
                logger.warning("Quote slide '%s' was not added, falling back to content", slide_spec.title)
            else:
                logger.warning("Quote slide '%s' missing text or author, falling back to content", slide_spec.title)
            self.add_content_slide(
                slide_spec.title, slide_spec.bullets, slide_spec.speaker_notes or "", image_path=slide_spec.image_path
            )
            return
        if slide_spec.layout == SlideLayout.STATISTICS:
            if slide_spec.statistics:
                stats_dicts = [{"value": stat.value, "label": stat.label} for stat in slide_spec.statistics]
                count_before = len(self.prs.slides)
                self.add_statistics_slide(slide_spec.title, stats_dicts, slide_spec.speaker_notes or "")
                if len(self.prs.slides) > count_before:
                    return
                logger.warning("Statistics slide '%s' data was insufficient, falling back to content", slide_spec.title)
            else:
                logger.warning("Statistics slide '%s' has no data, falling back to content", slide_spec.title)
            self.add_content_slide(
                slide_spec.title, slide_spec.bullets, slide_spec.speaker_notes or "", image_path=slide_spec.image_path
            )
            return
        if slide_spec.layout == SlideLayout.IMAGE:
            if slide_spec.image_path:
                count_before = len(self.prs.slides)
                self.add_fullscreen_image_slide(
                    slide_spec.image_path,
                    caption=slide_spec.image_caption or "",
                    overlay_title=slide_spec.title,
                    notes=slide_spec.speaker_notes or "",
                )
                if len(self.prs.slides) > count_before:
                    return
                logger.warning("Image slide '%s' was not added, falling back to content", slide_spec.title)
            else:
                logger.warning("Image slide '%s' has no image, falling back to content", slide_spec.title)
            self.add_content_slide(
                slide_spec.title, slide_spec.bullets, slide_spec.speaker_notes or ""
            )
            return
        if slide_spec.layout == SlideLayout.CHART:
            if slide_spec.chart_data:
                count_before = len(self.prs.slides)
                self.add_chart_slide(slide_spec.title, slide_spec.chart_data, slide_spec.speaker_notes or "")
                if len(self.prs.slides) > count_before:
                    return
                logger.warning("Chart slide '%s' data was insufficient, falling back to content", slide_spec.title)
            else:
                logger.warning("Chart slide '%s' has no data, falling back to content", slide_spec.title)
            self.add_content_slide(
                slide_spec.title, slide_spec.bullets, slide_spec.speaker_notes or "", image_path=slide_spec.image_path
            )
            return
        if slide_spec.layout == SlideLayout.CITATIONS:
            self.add_citations_slide(slide_spec.citations)
            return

        self.add_content_slide(
            slide_spec.title,
            slide_spec.bullets,
            slide_spec.speaker_notes or "",
            image_path=slide_spec.image_path,
        )

    def _add_blank_slide(self):
        blank_layout = None
        for layout in self.prs.slide_layouts:
            if layout.name.lower() == "blank":
                blank_layout = layout
                break
        if blank_layout is None:
            blank_layout = self.prs.slide_layouts[min(6, len(self.prs.slide_layouts) - 1)]
        slide = self.prs.slides.add_slide(blank_layout)
        self._apply_background(slide)
        return slide

    def _theme(self, key: str) -> Any:
        return self.current_style[key]

    def _slide_width_inches(self) -> float:
        w = self.prs.slide_width
        return (w if w is not None else 9144000) / 914400

    def _slide_height_inches(self) -> float:
        h = self.prs.slide_height
        return (h if h is not None else 6858000) / 914400

    def _apply_background(self, slide) -> None:
        if self.current_style.get("gradient"):
            try:
                bg = slide.background
                bg.fill.gradient()
                bg.fill.gradient_angle = 270
                bg.fill.gradient_stops[0].color.rgb = self._theme("bg_color")
                bg.fill.gradient_stops[1].color.rgb = self._theme("gradient_end")
                return
            except Exception as exc:
                logger.warning("Gradient fallback: %s", exc)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = self._theme("bg_color")

    def _set_notes(self, slide, notes: str) -> None:
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    def _add_title_text(self, slide, title: str, subtitle: str = "", eyebrow: str = "", hero: bool = False) -> None:
        left = self._theme("title_left_in")
        top = self._theme("title_top_in")
        width = self._theme("title_width_in")
        size = self._theme("hero_title_size_pt") if hero else self._theme("title_size_pt")
        current_top = top

        if eyebrow:
            eyebrow_box = slide.shapes.add_textbox(Inches(left), Inches(current_top), Inches(width), Inches(0.32))
            self._write_paragraphs(
                eyebrow_box.text_frame,
                [eyebrow],
                font_size=self._theme("eyebrow_size_pt"),
                color=self._theme("accent_color"),
                bold=True,
                font_name=self._theme("accent_font_name"),
                uppercase=True,
            )
            current_top += 0.34

        title_box = slide.shapes.add_textbox(Inches(left), Inches(current_top), Inches(width), Inches(1.8))
        self._write_paragraphs(
            title_box.text_frame,
            [title],
            font_size=size,
            color=self._theme("title_color"),
            bold=True,
            font_name=self._theme("title_font_name"),
            paragraph_spacing=0.05,
        )

        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Inches(left),
                Inches(current_top + 0.82 + self._theme("subtitle_gap_in")),
                Inches(width),
                Inches(0.9),
            )
            self._write_paragraphs(
                subtitle_box.text_frame,
                [subtitle],
                font_size=self._theme("subtitle_size_pt"),
                color=self._theme("subtitle_color"),
                font_name=self._theme("font_name"),
            )

        self._add_decoration_line(slide, current_top + 0.7 + self._theme("decoration_line_offset_in"))

    def _add_decoration_line(self, slide, y_position: float) -> None:
        if not self.current_style.get("decoration_line"):
            return
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(self._theme("title_left_in")),
            Inches(y_position),
            Inches(self._theme("decoration_line_width_in")),
            Inches(self._theme("decoration_line_height_in")),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._theme("accent_color")
        shape.line.fill.background()

    def _add_panel(self, slide, left: float, top: float, width: float, height: float, transparency: float | None = None):
        panel = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        panel.fill.solid()
        panel.fill.fore_color.rgb = self._theme("panel_fill_color")
        panel.fill.transparency = transparency if transparency is not None else self._theme("panel_fill_transparency")
        panel.line.color.rgb = self._theme("panel_border_color")
        panel.line.width = Pt(self._theme("panel_border_width"))
        return panel

    def _write_paragraphs(
        self,
        text_frame,
        paragraphs: list[str],
        *,
        font_size: float,
        color: RGBColor,
        font_name: str,
        bold: bool = False,
        italic: bool = False,
        alignment: PP_ALIGN | None = None,
        paragraph_spacing: float = 0.12,
        uppercase: bool = False,
        bullet: bool = False,
    ) -> None:
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(self._theme("content_padding_in"))
        text_frame.margin_right = Inches(self._theme("content_padding_in"))
        text_frame.margin_top = Inches(self._theme("content_padding_in"))
        text_frame.margin_bottom = Inches(self._theme("content_padding_in"))
        for index, text in enumerate(paragraphs):
            p = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            p.text = text.upper() if uppercase else text
            p.font.size = Pt(font_size)
            p.font.bold = bold
            p.font.italic = italic
            p.font.color.rgb = color
            p.font.name = font_name
            if alignment is not None:
                p.alignment = alignment
            p.space_after = Pt(font_size * paragraph_spacing)
            if bullet:
                p.level = 0

    def _cover_image(self, image_path: str, target_ratio: float) -> str:
        target_ratio = max(target_ratio, 0.001)
        try:
            raw_image = Image.open(image_path)
        except Image.DecompressionBombError as exc:
            raise RenderError("_cover_image", f"Image decompression bomb detected: {exc}") from exc
        with raw_image:
            w, h = raw_image.size
            if w * h > 25_000_000:
                raise RenderError("_cover_image", f"Image too large: {w}x{h} pixels")
            img = raw_image.convert("RGB")
            width, height = img.size
            current_ratio = width / max(height, 1)
            if current_ratio > target_ratio:
                target_width = max(int(height * target_ratio), 1)
                left = max((width - target_width) // 2, 0)
                img = img.crop((left, 0, left + target_width, height))
            else:
                target_height = max(int(width / target_ratio), 1)
                top = max((height - target_height) // 2, 0)
                img = img.crop((0, top, width, top + target_height))
            with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as temp:
                temp_path = temp.name
                try:
                    img.save(temp, format="PNG")
                except Exception:
                    try:
                        os.unlink(temp_path)
                    except OSError:
                        pass
                    raise
            return temp_path

    def _add_cover_picture(self, slide, image_path: str, left: float, top: float, width: float, height: float) -> bool:
        if not image_path or not os.path.exists(image_path):
            return False
        try:
            cropped_path = self._cover_image(image_path, width / max(height, 0.01))
        except (RenderError, OSError) as exc:
            logger.warning("Failed to process cover image %s: %s", image_path, exc)
            return False
        try:
            slide.shapes.add_picture(cropped_path, Inches(left), Inches(top), Inches(width), Inches(height))
            return True
        finally:
            try:
                os.unlink(cropped_path)
            except OSError:
                pass

    def add_title_slide(self, title: str, subtitle: str = "") -> None:
        slide = self._add_blank_slide()
        self._add_title_text(slide, title, subtitle, eyebrow="AutoPPT", hero=True)

        left = self._theme("title_left_in")
        top = self._theme("surface_top_in") + 1.0
        width = min(4.2, self._slide_width_inches() - left * 2)
        height = 1.45
        self._add_panel(slide, left, top, width, height, transparency=0.12)

        supporting = [
            "PowerPoint-native output",
            "Layout-aware slide planning",
            "Template and remix ready",
        ]
        support_box = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 0.08), Inches(width - 0.2), Inches(height - 0.16))
        self._write_paragraphs(
            support_box.text_frame,
            supporting,
            font_size=self._theme("bullet_size_pt"),
            color=self._theme("text_color"),
            font_name=self._theme("font_name"),
            paragraph_spacing=0.08,
            bullet=True,
        )

    def add_section_header(self, title: str) -> None:
        slide = self._add_blank_slide()
        self._add_title_text(slide, title, eyebrow="Section", hero=True)
        accent = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(self._theme("title_left_in")),
            Inches(4.55),
            Inches(2.4),
            Inches(0.24),
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = self._theme("accent_color")
        accent.line.fill.background()

    def add_content_slide(self, title: str, bullets: list[str], notes: str = "", image_path: str | None = None) -> None:
        slide = self._add_blank_slide()
        self._add_title_text(slide, title)

        top = self._theme("surface_top_in")
        height = self._theme("surface_height_in")
        left = self._theme("content_left_in")
        slide_width = self._slide_width_inches()
        outer_right = self._theme("image_inset_in")

        if image_path and os.path.exists(image_path):
            gap = self._theme("column_gap_in")
            image_width = self._theme("image_width_in")
            text_width = slide_width - left - outer_right - image_width - gap
            self._add_panel(slide, left, top, text_width, height)
            self._add_panel(slide, left + text_width + gap, top, image_width, height, transparency=0.0)
            text_box = slide.shapes.add_textbox(
                Inches(left + 0.05),
                Inches(top + 0.05),
                Inches(text_width - 0.1),
                Inches(height - 0.1),
            )
            self._write_paragraphs(
                text_box.text_frame,
                bullets or ["Content pending"],
                font_size=self._theme("bullet_size_pt"),
                color=self._theme("text_color"),
                font_name=self._theme("font_name"),
                paragraph_spacing=0.1,
                bullet=True,
            )
            self._add_cover_picture(slide, image_path, left + text_width + gap, top, image_width, height)
        else:
            width = slide_width - left - outer_right
            self._add_panel(slide, left, top, width, height)
            text_box = slide.shapes.add_textbox(
                Inches(left + 0.08),
                Inches(top + 0.08),
                Inches(width - 0.16),
                Inches(height - 0.16),
            )
            self._write_paragraphs(
                text_box.text_frame,
                bullets or ["Content pending"],
                font_size=self._theme("bullet_size_pt"),
                color=self._theme("text_color"),
                font_name=self._theme("font_name"),
                paragraph_spacing=0.1,
                bullet=True,
            )

        self._set_notes(slide, notes)

    def add_two_column_slide(
        self,
        title: str,
        left_bullets: list[str],
        right_bullets: list[str],
        left_title: str = "",
        right_title: str = "",
        notes: str = "",
    ) -> None:
        slide = self._add_blank_slide()
        self._add_title_text(slide, title)

        left = self._theme("content_left_in")
        top = self._theme("surface_top_in")
        gap = self._theme("column_gap_in")
        width = (self._slide_width_inches() - left - self._theme("image_inset_in") - gap) / 2
        height = self._theme("surface_height_in")

        self._add_panel(slide, left, top, width, height)
        self._add_panel(slide, left + width + gap, top, width, height)

        if left_title:
            left_head = slide.shapes.add_textbox(Inches(left + 0.08), Inches(top + 0.08), Inches(width - 0.16), Inches(0.45))
            self._write_paragraphs(
                left_head.text_frame,
                [left_title],
                font_size=self._theme("subtitle_size_pt"),
                color=self._theme("accent_color"),
                font_name=self._theme("accent_font_name"),
                bold=True,
            )
        if right_title:
            right_head = slide.shapes.add_textbox(Inches(left + width + gap + 0.08), Inches(top + 0.08), Inches(width - 0.16), Inches(0.45))
            self._write_paragraphs(
                right_head.text_frame,
                [right_title],
                font_size=self._theme("subtitle_size_pt"),
                color=self._theme("accent_color"),
                font_name=self._theme("accent_font_name"),
                bold=True,
            )

        left_box = slide.shapes.add_textbox(Inches(left + 0.08), Inches(top + 0.58), Inches(width - 0.16), Inches(height - 0.66))
        right_box = slide.shapes.add_textbox(
            Inches(left + width + gap + 0.08),
            Inches(top + 0.58),
            Inches(width - 0.16),
            Inches(height - 0.66),
        )
        self._write_paragraphs(
            left_box.text_frame,
            left_bullets or ["Pending"],
            font_size=self._theme("bullet_size_pt"),
            color=self._theme("text_color"),
            font_name=self._theme("font_name"),
            paragraph_spacing=0.1,
            bullet=True,
        )
        self._write_paragraphs(
            right_box.text_frame,
            right_bullets or ["Pending"],
            font_size=self._theme("bullet_size_pt"),
            color=self._theme("text_color"),
            font_name=self._theme("font_name"),
            paragraph_spacing=0.1,
            bullet=True,
        )
        self._set_notes(slide, notes)

    def add_comparison_slide(self, title: str, item_a: dict[str, Any], item_b: dict[str, Any], notes: str = "") -> None:
        self.add_two_column_slide(
            title=title,
            left_bullets=item_a.get("points", []),
            right_bullets=item_b.get("points", []),
            left_title=item_a.get("name", "Option A"),
            right_title=item_b.get("name", "Option B"),
            notes=notes,
        )

    def add_quote_slide(self, quote: str, author: str, context: str = "", notes: str = "") -> None:
        if not quote or not author:
            logger.warning("Skipping quote slide: missing quote or author")
            return
        slide = self._add_blank_slide()
        self._add_panel(slide, 1.0, 1.2, self._slide_width_inches() - 2.0, 4.8, transparency=0.12)

        quote_mark = slide.shapes.add_textbox(Inches(1.25), Inches(1.45), Inches(1.0), Inches(1.0))
        self._write_paragraphs(
            quote_mark.text_frame,
            ["\u201c"],
            font_size=self._theme("quote_mark_size_pt"),
            color=self._theme("accent_color"),
            font_name=self._theme("title_font_name"),
            bold=True,
        )

        quote_box = slide.shapes.add_textbox(Inches(2.0), Inches(2.0), Inches(self._slide_width_inches() - 3.1), Inches(2.5))
        self._write_paragraphs(
            quote_box.text_frame,
            [quote],
            font_size=self._theme("quote_size_pt"),
            color=self._theme("title_color"),
            font_name=self._theme("title_font_name"),
            italic=True,
            paragraph_spacing=0.06,
        )

        author_line = f"— {author}" + (f", {context}" if context else "")
        author_box = slide.shapes.add_textbox(Inches(2.0), Inches(5.0), Inches(self._slide_width_inches() - 3.1), Inches(0.5))
        self._write_paragraphs(
            author_box.text_frame,
            [author_line],
            font_size=self._theme("subtitle_size_pt"),
            color=self._theme("subtitle_color"),
            font_name=self._theme("font_name"),
        )
        self._set_notes(slide, notes)

    def add_chart_slide(self, title: str, chart_data: "ChartData", notes: str = "") -> None:
        from .data_types import ChartType

        if not chart_data.categories:
            logger.warning("Chart slide '%s' has no categories, skipping", title)
            return

        categories = list(chart_data.categories)
        values = list(chart_data.values)
        if len(values) != len(categories):
            shorter = min(len(values), len(categories))
            logger.warning(
                "Chart slide '%s': values length (%d) != categories length (%d), truncating to %d",
                title, len(values), len(categories), shorter,
            )
            categories = categories[:shorter]
            values = values[:shorter]

        if not categories:
            logger.warning("Chart slide '%s' has no data after truncation, skipping", title)
            return

        slide = self._add_blank_slide()
        self._add_title_text(slide, title)
        self._add_panel(
            slide,
            self._theme("chart_left_in") - 0.08,
            self._theme("chart_top_in") - 0.08,
            self._theme("chart_width_in") + 0.16,
            self._theme("chart_height_in") + 0.16,
            transparency=0.05,
        )

        chart_data_obj = CategoryChartData()
        chart_data_obj.categories = categories
        chart_data_obj.add_series(chart_data.series_name, values)

        chart_type_map = {
            ChartType.BAR: XL_CHART_TYPE.BAR_CLUSTERED,
            ChartType.COLUMN: XL_CHART_TYPE.COLUMN_CLUSTERED,
            ChartType.LINE: XL_CHART_TYPE.LINE,
            ChartType.PIE: XL_CHART_TYPE.PIE,
        }
        chart = slide.shapes.add_chart(
            chart_type_map.get(chart_data.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED),
            Inches(self._theme("chart_left_in")),
            Inches(self._theme("chart_top_in")),
            Inches(self._theme("chart_width_in")),
            Inches(self._theme("chart_height_in")),
            chart_data_obj,
        ).chart
        if chart_data.title:
            chart.has_title = True
            chart.chart_title.text_frame.text = chart_data.title
        self._set_notes(slide, notes)

    def add_citations_slide(self, citations: list[str]) -> None:
        if not citations:
            return
        slide = self._add_blank_slide()
        self._add_title_text(slide, "References")
        left = self._theme("content_left_in")
        top = self._theme("surface_top_in")
        width = self._slide_width_inches() - left - self._theme("image_inset_in")
        height = self._theme("surface_height_in")
        self._add_panel(slide, left, top, width, height, transparency=0.04)
        body = slide.shapes.add_textbox(Inches(left + 0.08), Inches(top + 0.08), Inches(width - 0.16), Inches(height - 0.16))
        self._write_paragraphs(
            body.text_frame,
            citations,
            font_size=self._theme("small_text_size_pt"),
            color=self._theme("text_color"),
            font_name=self._theme("font_name"),
            paragraph_spacing=0.04,
            bullet=True,
        )

    def add_fullscreen_image_slide(self, image_path: str, caption: str = "", overlay_title: str = "", notes: str = "") -> None:
        if not image_path or not os.path.exists(image_path):
            logger.warning("Image not found: %s", image_path)
            return
        slide = self._add_blank_slide()
        self._add_cover_picture(slide, image_path, 0, 0, self._slide_width_inches(), self._slide_height_inches())

        overlay = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0),
            Inches(0),
            self.prs.slide_width,
            self.prs.slide_height,
        )
        overlay.fill.solid()
        overlay.fill.fore_color.rgb = self._theme("image_overlay_color")
        overlay.fill.transparency = 1.0 - self._theme("image_overlay_opacity")
        overlay.line.fill.background()

        if overlay_title:
            title_box = slide.shapes.add_textbox(Inches(0.7), Inches(4.8), Inches(self._slide_width_inches() - 1.4), Inches(1.2))
            self._write_paragraphs(
                title_box.text_frame,
                [overlay_title],
                font_size=self._theme("hero_title_size_pt"),
                color=RGBColor(255, 255, 255),
                font_name=self._theme("title_font_name"),
                bold=True,
            )

        if caption:
            caption_box = slide.shapes.add_textbox(Inches(0.7), Inches(self._slide_height_inches() - 0.78), Inches(self._slide_width_inches() - 1.4), Inches(0.3))
            self._write_paragraphs(
                caption_box.text_frame,
                [caption],
                font_size=self._theme("caption_size_pt"),
                color=self._theme("image_caption_color"),
                font_name=self._theme("font_name"),
            )
        self._set_notes(slide, notes)

    def add_statistics_slide(self, title: str, stats: list[dict[str, str]], notes: str = "") -> None:
        if len(stats) > 4:
            logger.warning("Truncating statistics from %d to 4 items", len(stats))
            stats = stats[:4]
        if not stats:
            return
        slide = self._add_blank_slide()
        self._add_title_text(slide, title)

        left = self._theme("content_left_in")
        gap = self._theme("stat_card_gap_in")
        width = (self._slide_width_inches() - left - self._theme("image_inset_in") - gap * (len(stats) - 1)) / len(stats)
        width = max(width, 0.5)
        top = 2.2

        for index, stat in enumerate(stats):
            x = left + index * (width + gap)
            self._add_panel(slide, x, top, width, self._theme("stat_card_height_in"), transparency=0.04)

            value_box = slide.shapes.add_textbox(Inches(x + 0.08), Inches(top + 0.32), Inches(width - 0.16), Inches(0.8))
            self._write_paragraphs(
                value_box.text_frame,
                [stat.get("value", "N/A")],
                font_size=self._theme("stat_value_size_pt"),
                color=self._theme("accent_color"),
                font_name=self._theme("title_font_name"),
                bold=True,
                alignment=PP_ALIGN.CENTER,
            )

            label_box = slide.shapes.add_textbox(Inches(x + 0.08), Inches(top + 1.22), Inches(width - 0.16), Inches(0.7))
            self._write_paragraphs(
                label_box.text_frame,
                [stat.get("label", "")],
                font_size=self._theme("stat_label_size_pt"),
                color=self._theme("text_color"),
                font_name=self._theme("font_name"),
                alignment=PP_ALIGN.CENTER,
            )
        self._set_notes(slide, notes)

    def save(self, output_path: str) -> None:
        """Save the presentation, rejecting writes to sensitive system paths."""
        if ".." in output_path.replace("\\", "/").split("/"):
            raise RenderError("save", f"Path traversal detected: {output_path}")
        resolved = os.path.realpath(output_path)
        for prefix in Config.BLOCKED_SYSTEM_PREFIXES:
            if resolved.startswith(prefix):
                raise RenderError("save", f"Access to system path is not allowed: {output_path}")
        for segment in Config.BLOCKED_PATH_SEGMENTS:
            if segment in resolved:
                raise RenderError("save", f"Access to sensitive path is not allowed: {output_path}")
        if os.path.islink(output_path) or (os.path.exists(resolved) and os.path.islink(resolved)):
            raise RenderError("save", f"Refusing to write through symlink: {output_path}")
        try:
            self.prs.save(resolved)
        except Exception as exc:
            raise RenderError("save", str(exc)) from exc
        logger.info("Saved presentation to %s", resolved)
