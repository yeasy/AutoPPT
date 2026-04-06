"""
Template handler for working with existing PowerPoint presentations.
"""
import logging
import zipfile
from typing import Any
from pathlib import Path

from pptx import Presentation
try:
    from markitdown import MarkItDown
    HAS_MARKITDOWN = True
except ImportError:
    HAS_MARKITDOWN = False

from .config import Config

logger = logging.getLogger(__name__)

class TemplateHandler:
    """Handles loading and analysis of PowerPoint templates."""

    def __init__(self, template_path: str):
        """
        Initialize with a path to a PPTX template.
        """
        if ".." in str(template_path).replace("\\", "/").split("/"):
            raise ValueError(f"Path traversal detected: {template_path}")
        self.template_path = Path(template_path).resolve()
        resolved_str = str(self.template_path)
        for prefix in Config.BLOCKED_SYSTEM_PREFIXES:
            if resolved_str.startswith(prefix):
                raise ValueError(f"Access to system path is not allowed: {template_path}")
        if not self.template_path.exists():
            raise FileNotFoundError(f"Template not found: {template_path}")
        if self.template_path.suffix.lower() != ".pptx":
            raise ValueError(f"Invalid template file type: {self.template_path.suffix} (expected .pptx)")
        file_size = self.template_path.stat().st_size
        if file_size > Config.MAX_TEMPLATE_BYTES:
            raise ValueError(f"Template file too large: {file_size} bytes (max {Config.MAX_TEMPLATE_BYTES})")

        try:
            with zipfile.ZipFile(str(self.template_path), "r") as zf:
                total = sum(info.file_size for info in zf.infolist())
                if total > Config.MAX_DECOMPRESSED_BYTES:
                    raise ValueError(
                        f"Template decompressed size ({total} bytes) exceeds limit ({Config.MAX_DECOMPRESSED_BYTES})"
                    )
        except zipfile.BadZipFile as exc:
            raise ValueError(f"Invalid PPTX file: {exc}") from exc

        self.prs = Presentation(str(self.template_path))
        self.layouts = self._analyze_layouts()

    def _analyze_layouts(self) -> dict[int, dict[str, Any]]:
        """Analyze available slide layouts."""
        layouts: dict[int, dict[str, Any]] = {}
        flat_index = 0

        for i, slide_master in enumerate(self.prs.slide_masters):
            for j, layout in enumerate(slide_master.slide_layouts):
                layout_info = {
                    "index": j,
                    "master_index": i,
                    "name": layout.name,
                    "placeholders": [
                        {
                            "idx": ph.placeholder_format.idx,
                            "type": str(ph.placeholder_format.type),
                            "name": ph.name
                        }
                        for ph in layout.placeholders
                    ]
                }
                layouts[flat_index] = layout_info
                flat_index += 1

        return layouts

    def extract_text_content(self) -> str:
        """Extract text content from the template using MarkItDown."""
        if not HAS_MARKITDOWN:
            logger.warning("markitdown not installed, cannot extract text content")
            return ""

        try:
            md = MarkItDown()
            result = md.convert(str(self.template_path))
            return str(result.text_content)
        except Exception as exc:
            logger.error("Failed to extract text from template: %s", exc)
            return ""

    def get_layout_by_name(self, name_pattern: str) -> int | None:
        """Find a layout index by name matching."""
        name_pattern = name_pattern.lower()
        for idx, info in self.layouts.items():
            if name_pattern in info["name"].lower():
                return idx
        return None

    def get_best_layout_for_type(self, slide_type: str) -> int:
        """
        Get the best layout index for a given slide type.
        """
        # Common PowerPoint layout names
        mappings = {
            "title": ["title slide", "intro"],
            "section": ["section header", "segue"],
            "content": ["title and content", "content"],
            "two_column": ["two content", "comparison"],
            "comparison": ["two content", "comparison"],
            "quote": ["title and content", "content"],
            "chart": ["title and content", "content"],
            "statistics": ["title and content", "content"],
            "citations": ["title and content", "content"],
            "blank": ["blank"],
            "picture": ["picture", "image"],
            "image": ["picture", "image"],
        }

        patterns = mappings.get(slide_type, [])
        for pattern in patterns:
            idx = self.get_layout_by_name(pattern)
            if idx is not None:
                return idx

        # Fallback to first non-title layout for content or unknown types
        return 1 if 1 in self.layouts else 0
