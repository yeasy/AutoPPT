from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from .data_types import UserPresentation, SlideConfig

class PPTRenderer:
    def __init__(self, template_path: str = None):
        if template_path:
            self.prs = Presentation(template_path)
        else:
            self.prs = Presentation()
            
    def add_title_slide(self, title: str, subtitle: str = ""):
        """Add a title slide with styled colors."""
        slide_layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._apply_background(slide)
        
        title_shape = slide.shapes.title
        subtitle_shape = slide.placeholders[1]
        
        title_shape.text = title
        self._style_text_shape(title_shape, is_title=True)
        
        subtitle_shape.text = subtitle
        self._style_text_shape(subtitle_shape, is_title=False)
        
    def add_section_header(self, title: str):
        """Add a section header slide."""
        slide_layout = self.prs.slide_layouts[2] # Usually Section Header
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._apply_background(slide)
        
        if slide.shapes.title:
            slide.shapes.title.text = title
            self._style_text_shape(slide.shapes.title, is_title=True)
            
    def apply_style(self, style_name: str):
        """Setup global presentation style (colors, fonts)."""
        styles = {
            "technology": {
                "title_color": RGBColor(0, 102, 204),
                "text_color": RGBColor(200, 200, 255),
                "font_name": "Arial",
                "bg_color": RGBColor(10, 10, 40) 
            },
            "nature": {
                "title_color": RGBColor(34, 139, 34),
                "text_color": RGBColor(50, 80, 50),
                "font_name": "Georgia",
                "bg_color": RGBColor(245, 255, 250)
            },
            "creative": {
                "title_color": RGBColor(200, 50, 150),
                "text_color": RGBColor(60, 40, 60),
                "font_name": "Verdana",
                "bg_color": RGBColor(255, 250, 240)
            },
            "minimalist": {
                "title_color": RGBColor(40, 40, 40),
                "text_color": RGBColor(80, 80, 80),
                "font_name": "Arial",
                "bg_color": RGBColor(255, 255, 255)
            }
        }
        self.current_style = styles.get(style_name.lower(), styles["minimalist"])

    def _apply_background(self, slide):
        """Apply the current style's background color."""
        if hasattr(self, 'current_style'):
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = self.current_style["bg_color"]

    def _style_text_shape(self, shape, is_title=False):
        """Apply colors and fonts to a shape's text."""
        if not hasattr(self, 'current_style') or not shape.has_text_frame:
            return
            
        color = self.current_style["title_color"] if is_title else self.current_style["text_color"]
        
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = self.current_style["font_name"]
                run.font.color.rgb = color

    def add_content_slide(self, title: str, bullets: list, notes: str = "", image_path: str = None):
        slide_layout = self.prs.slide_layouts[1] # 1 is usually Title and Content
        slide = self.prs.slides.add_slide(slide_layout)
        
        self._apply_background(slide)
        
        # Title
        title_shape = slide.shapes.title
        title_shape.text = title
        self._style_text_shape(title_shape, is_title=True)
        
        # Content (Bullets)
        body_shape = slide.placeholders[1]
        
        # If image, we resize the body shape to make room
        if image_path and os.path.exists(image_path):
            body_shape.width = Inches(5.5)
            try:
                slide.shapes.add_picture(image_path, Inches(6), Inches(1.5), height=Inches(5))
            except Exception as e:
                print(f"Warning: Failed to add image {image_path}: {e}")
        
        tf = body_shape.text_frame
        tf.word_wrap = True
        tf.clear() 
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            
        self._style_text_shape(body_shape, is_title=False)
            
        # Add Notes
        if notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = notes
            
    def add_citations_slide(self, citations: list):
        if not citations:
            return
            
        slide_layout = self.prs.slide_layouts[1]
        slide = self.prs.slides.add_slide(slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = "References"
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for cit in citations:
            p = tf.add_paragraph()
            p.text = cit
            p.font.size = Pt(12)
            
    def save(self, output_path: str):
        self.prs.save(output_path)
